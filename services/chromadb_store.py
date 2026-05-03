import chromadb
from chromadb.utils import embedding_functions

chroma_client = chromadb.PersistentClient(path="./chroma_db")
ef = embedding_functions.DefaultEmbeddingFunction()

policy_collection = chroma_client.get_or_create_collection(
    name="company_policies", embedding_function=ef
)
tool_knowledge_collection = chroma_client.get_or_create_collection(
    name="tool_knowledge", embedding_function=ef
)


def _chunk_text(text: str, chunk_size: int = 400, overlap: int = 60) -> list:
    words  = text.split()
    chunks = []
    start  = 0
    while start < len(words):
        end = start + chunk_size
        chunks.append(" ".join(words[start:end]))
        start += chunk_size - overlap
    return [c.strip() for c in chunks if len(c.strip()) > 60]


def _classify_chunk_to_tool(chunk: str, tool_names: list) -> str:
    from services.llm_client import HAS_AZURE, _azure_client, call_llm
    if not HAS_AZURE or not _azure_client or not tool_names:
        return "unclassified"

    names_list = "\n".join(f"- {n}" for n in tool_names)
    system_msg = (
        "You are a document classifier. Given a text chunk from a tool document, "
        "identify which AI tool it describes from the provided list. "
        "Reply with ONLY the exact tool name from the list, or 'unclassified' if unsure."
    )
    user_msg = (
        f"Tool names to choose from:\n{names_list}\n\n"
        f"Text chunk:\n{chunk[:800]}\n\n"
        "Which tool does this text describe? Reply with only the exact tool name or 'unclassified'."
    )
    try:
        result = call_llm(system_msg, user_msg, max_tokens=30, temperature=0.0)
        result = result.strip().strip('"').strip("'")
        if result in tool_names:
            return result
        result_lower = result.lower()
        for name in tool_names:
            if name.lower() == result_lower or name.lower() in result_lower:
                return name
        return "unclassified"
    except Exception:
        return "unclassified"


def _extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    from services.file_utils import extract_pdf_text, extract_docx_text
    fname = filename.lower()
    text  = ""
    try:
        if fname.endswith(".pdf"):
            text = extract_pdf_text(file_bytes)
        elif fname.endswith(".docx"):
            text = extract_docx_text(file_bytes)
        elif fname.endswith(".pptx"):
            try:
                from pptx import Presentation
                from io import BytesIO
                prs = Presentation(BytesIO(file_bytes))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text += shape.text.strip() + "\n"
            except Exception:
                text = file_bytes.decode("utf-8", errors="ignore")
        else:
            text = file_bytes.decode("utf-8", errors="ignore")
    except Exception:
        text = ""
    return text


def ingest_tool_document_direct(file_bytes: bytes, filename: str, tool_name: str) -> dict:
    text = _extract_text_from_file(file_bytes, filename)
    if not text.strip():
        return {}

    chunks = _chunk_text(text)
    count  = 0

    for idx, chunk in enumerate(chunks):
        chunk_id = f"{tool_name}__{filename}__chunk_{idx}"
        try:
            tool_knowledge_collection.upsert(
                ids       = [chunk_id],
                documents = [chunk],
                metadatas = [{
                    "tool_name":   tool_name,
                    "source_file": filename,
                    "chunk_index": idx,
                }],
            )
            count += 1
        except Exception:
            pass

    return {tool_name: count}


def ingest_tool_document(file_bytes: bytes, filename: str, tool_names: list) -> dict:
    text = _extract_text_from_file(file_bytes, filename)
    if not text.strip():
        return {}

    chunks  = _chunk_text(text)
    summary = {}

    for idx, chunk in enumerate(chunks):
        tool_name = _classify_chunk_to_tool(chunk, tool_names)
        chunk_id  = f"{filename}__chunk_{idx}"

        try:
            tool_knowledge_collection.upsert(
                ids        = [chunk_id],
                documents  = [chunk],
                metadatas  = [{
                    "tool_name":   tool_name,
                    "source_file": filename,
                    "chunk_index": idx,
                }],
            )
        except Exception:
            pass

        summary[tool_name] = summary.get(tool_name, 0) + 1

    return summary


def query_tool_knowledge(user_input: str, tool_names: list, n_results: int = 5) -> dict:
    if not tool_names:
        return {}
    try:
        results = tool_knowledge_collection.query(
            query_texts = [user_input],
            n_results   = min(n_results, 10),
        )
    except Exception:
        return {}

    docs      = results.get("documents",  [[]])[0]
    metas     = results.get("metadatas",  [[]])[0]
    knowledge = {}

    for doc, meta in zip(docs, metas):
        t = meta.get("tool_name", "unclassified")
        if t in tool_names:
            knowledge.setdefault(t, []).append(doc)

    return knowledge


def get_tool_knowledge_status() -> list:
    try:
        all_items = tool_knowledge_collection.get(include=["metadatas"])
        metas     = all_items.get("metadatas", [])
    except Exception:
        return []

    status = {}
    for m in metas:
        t     = m.get("tool_name", "unclassified")
        fname = m.get("source_file", "")
        if t not in status:
            status[t] = {"tool_name": t, "chunk_count": 0, "source_files": set()}
        status[t]["chunk_count"]  += 1
        status[t]["source_files"].add(fname)

    return [
        {**v, "source_files": sorted(v["source_files"])}
        for v in sorted(status.values(), key=lambda x: x["chunk_count"], reverse=True)
    ]
