import json
import uuid
import io
import os
from datetime import datetime
from typing import List
from schemas import RunRequest, FeedbackRequest, PromptVersionRequest, RefinementRequest, AuditUpdateRequest, ClarifyRequest, ClarifyAnswerRequest, ChatGatherRequest, ChatSummarizeRequest, RegisterToolRequest, RegisterScenarioRequest
import auth as _auth


from fastapi import APIRouter, UploadFile, File, HTTPException, Form
from service import (
    orchestrator,
    get_db,
    policy_collection,
    tool_knowledge_collection,
    AI_TOOLS_REGISTRY,
    SYSTEM_VERSION,
    call_llm,
    call_llm_messages,
    ingest_tool_document,
    get_tool_knowledge_status,
    log_tool_change,
)

router = APIRouter()

# ══════════════════════════════════════════════════════════════════════════════
# AZURE BLOB STORAGE — FEEDBACK
# ══════════════════════════════════════════════════════════════════════════════
def _get_blob_container():
    from azure.storage.blob import BlobServiceClient
    account_name = os.getenv("ACCOUNT_NAME", "")
    account_key  = os.getenv("ACCOUNT_KEY", "")
    container    = os.getenv("AZURE_STORAGE_CONTAINER_NAME", "ai-navigator-feedback")
    conn_str = (
        f"DefaultEndpointsProtocol=https;"
        f"AccountName={account_name};"
        f"AccountKey={account_key};"
        f"EndpointSuffix=core.windows.net"
    )
    client = BlobServiceClient.from_connection_string(conn_str)
    return client.get_container_client(container)


def _content_settings(content_type: str):
    from azure.storage.blob import ContentSettings
    return ContentSettings(content_type=content_type)


# ══════════════════════════════════════════════════════════════════════════════
# FILE TEXT EXTRACTION HELPERS
# ══════════════════════════════════════════════════════════════════════════════
def _extract_pdf_text(content: bytes) -> str:
    """Extract text from a PDF file using pypdf."""
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(content))
        parts  = []
        for page in reader.pages:
            try:
                parts.append(page.extract_text() or "")
            except Exception:
                pass
        return "\n".join(parts)
    except ImportError:
        # Fallback: try pdfplumber
        try:
            import pdfplumber
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                return "\n".join(
                    page.extract_text() or "" for page in pdf.pages
                )
        except ImportError:
            raise HTTPException(
                500,
                "PDF parsing library not installed. Run: pip install pypdf"
            )


def _extract_docx_text(content: bytes) -> str:
    """Extract text from a .docx file using python-docx."""
    try:
        from docx import Document
        doc   = Document(io.BytesIO(content))
        parts = [para.text for para in doc.paragraphs if para.text.strip()]
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    parts.append(row_text)
        return "\n".join(parts)
    except ImportError:
        raise HTTPException(
            500,
            "python-docx library not installed. Run: pip install python-docx"
        )


# ══════════════════════════════════════════════════════════════════════════════
# ORCHESTRATOR
# ══════════════════════════════════════════════════════════════════════════════
@router.post("/api/run")
async def run_orchestrator(req: RunRequest):
    if not req.user_input.strip():
        raise HTTPException(400, "Input cannot be empty")

    result = orchestrator.invoke({
        "user_input":       req.user_input,
        "role":             req.role or "general",
        "task_type":        req.task_type or "general",
        "data_sensitivity": req.data_sensitivity or "general",
        "intent": "", "industry": "",
        "recommended_tool": "", "tool_reason": "",
        "tool_confidence": "", "tool_confidence_pct": 0,
        "tool_confidence_explanation": "",
        "tool_alternatives": [], "tool_alternative_reasons": [],
        "tool_alternative_confidence_pcts": [], "tool_alternative_urls": [],
        "policy_flags": [], "policies": [],
        "policy_summary": "", "policy_blocked": False,
        "internal_results": None, "external_results": None,
        "corlo_prompt": "", "prompt_version": "1.0",
        "llm_output": "", "token_estimate": 0, "error": None,
    })

    audit_id  = str(uuid.uuid4())
    tool_info = AI_TOOLS_REGISTRY.get(result["recommended_tool"], {})

    is_blocked     = result.get("policy_blocked", False)
    policy_summary = result.get("policy_summary", "")

    stored_role  = (req.role or "").strip() or "general"
    stored_email = (req.user_email or "").strip().lower()

    conn = get_db()
    conn.execute(
        """INSERT INTO audit_log
            (id, created_at, raw_input, intent, industry,
             recommended_tool, tool_reason, tool_confidence,
             policy_flags, retrieved_policies, final_prompt,
             prompt_version, model_used, output, token_estimate,
             system_version, policy_blocked, policy_summary, role, user_email)
           VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)""",
        (
            audit_id, datetime.utcnow().isoformat(),
            result["user_input"], result["intent"], result["industry"],
            result["recommended_tool"], result["tool_reason"],
            result["tool_confidence"], json.dumps(result["policy_flags"]),
            json.dumps(result["policies"]), result["corlo_prompt"],
            result.get("prompt_version", "1.0"),
            result["recommended_tool"], result["llm_output"],
            result["token_estimate"], SYSTEM_VERSION,
            1 if is_blocked else 0,
            policy_summary,
            stored_role,
            stored_email,
        ),
    )
    conn.commit()
    conn.close()

    return {
        "audit_id":          audit_id,
        "intent":            result["intent"],
        "industry":          result["industry"],
        "recommended_tool":  result["recommended_tool"],
        "tool_reason":       result["tool_reason"],
        "tool_confidence":   result["tool_confidence"],
        "tool_alternatives":                result["tool_alternatives"],
        "tool_alternative_reasons":         result.get("tool_alternative_reasons", []),
        "tool_alternative_confidence_pcts": result.get("tool_alternative_confidence_pcts", []),
        "tool_alternative_urls":            result.get("tool_alternative_urls", []),
        "tool_confidence_pct":              result.get("tool_confidence_pct", 0),
        "tool_confidence_explanation":      result.get("tool_confidence_explanation", ""),
        "tool_icon":         tool_info.get("icon", "🤖"),
        "tool_category":     tool_info.get("category", ""),
        "tool_url":          result.get("tool_url") or tool_info.get("url", ""),
        "tool_is_internal":  tool_info.get("is_internal", False),
        "policy_flags":      result["policy_flags"],
        "policies":          result["policies"],
        "policy_summary":    policy_summary,
        "policy_blocked":    is_blocked,
        "corlo_prompt":      result["corlo_prompt"],
        "prompt_version":    result.get("prompt_version", "1.0"),
        "output":            result["llm_output"],
        "token_estimate":    result["token_estimate"],
        "role":              req.role,
        "task_type":         req.task_type,
        "data_sensitivity":  req.data_sensitivity,
    }


# ══════════════════════════════════════════════════════════════════════════════
# CLARIFICATION  —  assess whether the task description is clear enough,
#                   and if not, return 1-3 short yes/no or one-word questions.
# ══════════════════════════════════════════════════════════════════════════════
@router.post("/api/clarify")
async def check_clarity(req: ClarifyRequest):
    """
    Evaluate whether the user's task description is specific enough for accurate
    tool recommendation. Returns:
      - needs_clarification: bool
      - questions: list[str]  (1-4 short, task-specific questions)
    """
    if not req.user_input.strip():
        raise HTTPException(400, "Input cannot be empty")

    system_msg = (
        "You are a helpful assistant trying to understand what a user wants to do "
        "so you can recommend the right AI tool for them. "
        "Return ONLY valid JSON — no markdown, no extra text."
    )

    user_msg = (
        f"A user described their task as:\n\"{req.user_input}\"\n\n"
        "Analyse the task carefully and ask ONLY the questions that are genuinely needed "
        "to recommend the right tool. Use your judgement:\n"
        "- If the task is simple and clear, ask 1 question or none at all.\n"
        "- If the task is moderately ambiguous, ask 2-3 questions.\n"
        "- Only ask more if there are truly that many distinct unclear parts.\n\n"
        "For each unclear part you find, write one short question asking the user to clarify "
        "that specific part. Only ask about things they actually wrote — never bring up new topics.\n\n"
        "If everything they wrote is already clear, return needs_clarification=false and no questions.\n\n"
        "Return ONLY this JSON:\n"
        "{\"needs_clarification\": true or false, \"questions\": []}"
    )

    try:
        raw = call_llm(system_msg, user_msg, max_tokens=300, temperature=0.2)
        raw = raw.replace("```json", "").replace("```", "").strip()
        data = json.loads(raw)
        needs = bool(data.get("needs_clarification", False))
        questions = [str(q).strip() for q in (data.get("questions") or []) if str(q).strip()]
        return {"needs_clarification": needs, "questions": questions}
    except Exception:
        return {"needs_clarification": False, "questions": []}


@router.post("/api/clarify-merge")
async def merge_clarification(req: ClarifyAnswerRequest):
    """
    Merge the user's original input with Q&A answers into an enriched description
    that the /api/run endpoint can use for accurate tool recommendation.
    """
    if not req.user_input.strip():
        raise HTTPException(400, "Input cannot be empty")

    qa_pairs = "\n".join(
        f"Q: {q}\nA: {a}"
        for q, a in zip(req.questions, req.answers)
        if q.strip() and a.strip()
    )

    system_msg = (
        "You are a task description enricher for an enterprise AI tool recommender. "
        "Combine the original task description with the user's clarification answers "
        "into a single, clear, specific task description. "
        "Output ONLY the enriched description — no preamble, no explanation, no JSON."
    )
    user_msg = (
        f"ORIGINAL TASK DESCRIPTION:\n{req.user_input}\n\n"
        f"CLARIFICATION Q&A:\n{qa_pairs}\n\n"
        "Write a single enriched task description that incorporates all the above details. "
        "Keep it natural and concise (2-4 sentences max)."
    )

    try:
        enriched = call_llm(system_msg, user_msg, max_tokens=300, temperature=0.1)
        return {"enriched_input": enriched.strip()}
    except Exception:
        combined = req.user_input
        if qa_pairs:
            combined += " " + " ".join(req.answers)
        return {"enriched_input": combined.strip()}


# ══════════════════════════════════════════════════════════════════════════════
# REFINEMENT  —  user adds a comment on the existing LLM output
# ══════════════════════════════════════════════════════════════════════════════
@router.post("/api/refine")
async def refine_output(req: RefinementRequest):
    """
    Accept a user comment and revise the CORLO prompt accordingly.
    The revised prompt is returned so the frontend can display it
    and the user can copy it into the recommended AI tool.
    """
    if not req.comment.strip():
        raise HTTPException(400, "Comment cannot be empty")

    system_msg = (
        f"You are an expert prompt engineer helping a {req.role} working on a "
        f"{req.task_type} task in the {req.industry} industry. "
        f"Data sensitivity: {req.data_sensitivity}. "
        f"Target tool: {req.recommended_tool}. "
        "Your job is to revise an existing CORLO prompt based on the user's feedback. "
        "The CORLO prompt is structured in 5 sections: ROLE, CONTEXT, OBJECTIVE, LIMITATIONS, OUTPUT. "
        "Apply the user's comment precisely — change only what they ask, keep everything else. "
        "Return only the complete revised CORLO prompt — no preamble, no explanation."
    )
    user_msg = (
        f"ORIGINAL USER REQUEST:\n{req.user_input}\n\n"
        f"CURRENT CORLO PROMPT (what you are revising):\n{req.corlo_prompt}\n\n"
        f"USER FEEDBACK / REVISION COMMENT:\n{req.comment}\n\n"
        "Instructions:\n"
        "- Read the comment carefully — it tells you exactly what to change in the prompt.\n"
        "- If they say 'make it shorter' → tighten the OBJECTIVE and OUTPUT sections.\n"
        "- If they say 'add X' → insert it in the most appropriate section.\n"
        "- If they say 'focus on Y' → adjust the OBJECTIVE and OUTPUT accordingly.\n"
        "- Keep the 5-section structure (ROLE, CONTEXT, OBJECTIVE, LIMITATIONS, OUTPUT).\n"
        "- Return the complete revised prompt only."
    )

    try:
        revised = call_llm(system_msg, user_msg, max_tokens=2000)
    except Exception as e:
        raise HTTPException(500, f"LLM refinement failed: {str(e)}")

    return {
        "audit_id":       req.audit_id,
        "revised_output": revised,
    }


# ══════════════════════════════════════════════════════════════════════════════
# FEEDBACK — stored in Azure Blob Storage
# ══════════════════════════════════════════════════════════════════════════════
@router.post("/api/feedback")
async def submit_feedback(
    email:      str  = Form(""),
    rating:     int  = Form(...),
    comment:    str  = Form(""),
    issue_type: str  = Form(""),
    audit_id:   str  = Form(""),
    source:     str  = Form("form"),
    files:      List[UploadFile] = File(default=[]),
):
    feedback_id  = str(uuid.uuid4())
    created_at   = datetime.utcnow().isoformat()
    folder       = f"feedback/{created_at[:10]}_{feedback_id}"

    metadata = {
        "id":         feedback_id,
        "audit_id":   audit_id,
        "email":      email,
        "rating":     rating,
        "comment":    comment,
        "issue_type": issue_type,
        "source":     source,
        "created_at": created_at,
        "files":      [],
    }

    try:
        container = _get_blob_container()

        uploaded_files = []
        for f in files:
            if not f.filename:
                continue
            raw = await f.read()
            if not raw:
                continue
            safe_name    = f.filename.replace(" ", "_")
            blob_name    = f"{folder}/attachments/{safe_name}"
            content_type = f.content_type or "application/octet-stream"
            container.upload_blob(
                name=blob_name,
                data=raw,
                overwrite=True,
                content_settings=_content_settings(content_type),
            )
            uploaded_files.append(safe_name)

        metadata["files"] = uploaded_files

        meta_blob = f"{folder}/metadata.json"
        container.upload_blob(
            name=meta_blob,
            data=json.dumps(metadata, indent=2).encode("utf-8"),
            overwrite=True,
            content_settings=_content_settings("application/json"),
        )
    except Exception as e:
        raise HTTPException(500, f"Blob upload failed: {str(e)}")

    return {"status": "ok", "feedback_id": feedback_id}


@router.get("/api/feedback-list")
async def list_feedback(page: int = 1, per_page: int = 20, rating: int = 0, search: str = ""):
    try:
        container = _get_blob_container()
        blobs = list(container.list_blobs(name_starts_with="feedback/"))
    except Exception as e:
        raise HTTPException(500, f"Blob list failed: {str(e)}")

    meta_blobs = [b for b in blobs if b.name.endswith("/metadata.json")]

    all_feedbacks = []
    for blob in meta_blobs:
        try:
            data = container.download_blob(blob.name).readall()
            entry = json.loads(data)
            all_feedbacks.append(entry)
        except Exception:
            continue

    all_feedbacks.sort(key=lambda x: x.get("created_at", ""), reverse=True)

    ratings_all = [f.get("rating", 0) for f in all_feedbacks if f.get("rating")]
    avg_rating  = round(sum(ratings_all) / len(ratings_all), 1) if ratings_all else None

    from collections import Counter
    dist_counter = Counter(f.get("rating") for f in all_feedbacks if f.get("rating"))
    distribution = [{"rating": r, "count": c} for r, c in sorted(dist_counter.items())]

    filtered = all_feedbacks
    if rating > 0:
        filtered = [f for f in filtered if f.get("rating") == rating]
    if search.strip():
        q = search.strip().lower()
        filtered = [
            f for f in filtered
            if q in (f.get("email") or "").lower()
            or q in (f.get("comment") or "").lower()
            or q in (f.get("issue_type") or "").lower()
        ]

    total = len(filtered)
    offset = (page - 1) * per_page
    page_items = filtered[offset: offset + per_page]

    return {
        "total":        total,
        "page":         page,
        "per_page":     per_page,
        "avg_rating":   avg_rating,
        "distribution": distribution,
        "feedbacks":    page_items,
    }


@router.get("/api/feedback-attachments/{feedback_id}")
async def get_feedback_attachments(feedback_id: str):
    try:
        container  = _get_blob_container()
        prefix     = f"feedback/"
        all_blobs  = list(container.list_blobs(name_starts_with=prefix))
        folder_blob = next(
            (b for b in all_blobs if feedback_id in b.name and b.name.endswith("/metadata.json")),
            None
        )
        if not folder_blob:
            raise HTTPException(404, "Feedback not found")

        folder = folder_blob.name.replace("/metadata.json", "")
        attach_prefix = f"{folder}/attachments/"
        attach_blobs  = [b for b in all_blobs if b.name.startswith(attach_prefix)]

        urls = []
        for blob in attach_blobs:
            from azure.storage.blob import generate_blob_sas, BlobSasPermissions
            sas = generate_blob_sas(
                account_name   = os.getenv("ACCOUNT_NAME", ""),
                container_name = os.getenv("AZURE_STORAGE_CONTAINER_NAME", "ai-navigator-feedback"),
                blob_name      = blob.name,
                account_key    = os.getenv("ACCOUNT_KEY", ""),
                permission     = BlobSasPermissions(read=True),
                expiry         = datetime.utcnow().replace(hour=23, minute=59, second=59),
            )
            account_name = os.getenv("ACCOUNT_NAME", "")
            container_name = os.getenv("AZURE_STORAGE_CONTAINER_NAME", "ai-navigator-feedback")
            url = f"https://{account_name}.blob.core.windows.net/{container_name}/{blob.name}?{sas}"
            urls.append({"name": blob.name.split("/")[-1], "url": url})

        return {"files": urls}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, str(e))


# ══════════════════════════════════════════════════════════════════════════════
# POLICIES
# ══════════════════════════════════════════════════════════════════════════════
@router.post("/api/upload-policy")
async def upload_policy(file: UploadFile = File(...)):
    content  = await file.read()
    filename = file.filename or ""
    ext      = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    # ── Extract text based on file type ────────────────────────────────────
    if ext == "pdf":
        text = _extract_pdf_text(content)
    elif ext == "docx":
        text = _extract_docx_text(content)
    else:
        # Plain text / markdown fallback
        text = content.decode("utf-8", errors="ignore")

    if not text or not text.strip():
        raise HTTPException(400, "Could not extract any text from the uploaded file. "
                                 "Ensure the file contains readable text.")

    chunks    = [text[i:i+500] for i in range(0, len(text), 500) if text[i:i+500].strip()]
    if not chunks:
        raise HTTPException(400, "File appears empty after text extraction.")

    ids       = [f"{filename}_{i}" for i in range(len(chunks))]
    metadatas = [{"source": filename, "chunk": i} for i in range(len(chunks))]

    try:
        existing = policy_collection.get(where={"source": filename})
        if existing["ids"]:
            policy_collection.delete(ids=existing["ids"])
    except Exception:
        pass

    policy_collection.add(documents=chunks, ids=ids, metadatas=metadatas)
    return {"status": "ok", "filename": filename, "chunks_indexed": len(chunks)}


@router.get("/api/policies")
async def list_policies():
    try:
        result  = policy_collection.get()
        sources = list(set(m.get("source", "unknown") for m in (result.get("metadatas") or [])))
        return {"sources": sources, "total_chunks": len(result.get("ids") or [])}
    except Exception:
        return {"sources": [], "total_chunks": 0}


@router.delete("/api/policies/{filename}")
async def delete_policy(filename: str):
    try:
        existing = policy_collection.get(where={"source": filename})
        if existing["ids"]:
            policy_collection.delete(ids=existing["ids"])
            return {"status": "ok"}
        return {"status": "not_found"}
    except Exception as e:
        raise HTTPException(500, str(e))


# ══════════════════════════════════════════════════════════════════════════════
# TOOLS REGISTRY
# ══════════════════════════════════════════════════════════════════════════════
@router.get("/api/tools")
async def get_tools():
    return {
        name: {
            "description": info.get("description") or "",
            "desc_content": (info.get("raw_data") or {}).get("Desc_Content", ""),
            "category":    info.get("category", ""),
            "icon":        info.get("icon", ""),
            "best_for":    info.get("best_for") or [],
            "not_for":     info.get("not_for") or [],
            "strong_signals": info.get("strong_signals") or [],
            "weak_signals":   info.get("weak_signals") or [],
            "roles":       info.get("roles") or [],
            "output_type": info.get("output_type", ""),
            "url":         info.get("url", ""),
            "is_internal": info.get("is_internal", False),
            "raw_data":    info.get("raw_data", {}),
            "_source":     info.get("_source", "registry"),
        }
        for name, info in AI_TOOLS_REGISTRY.items()
    }


# ══════════════════════════════════════════════════════════════════════════════
# AUDIT & ANALYTICS
# ══════════════════════════════════════════════════════════════════════════════
@router.get("/api/audit")
async def get_audit_log(limit: int = 20, user_email: str = ""):
    conn = get_db()
    if user_email.strip():
        rows = conn.execute(
            "SELECT * FROM audit_log WHERE LOWER(user_email) = ? ORDER BY created_at DESC LIMIT ?",
            (user_email.strip().lower(), limit)
        ).fetchall()
    else:
        rows = conn.execute(
            "SELECT * FROM audit_log ORDER BY created_at DESC LIMIT ?", (limit,)
        ).fetchall()
    conn.close()
    return [dict(r) for r in rows]

@router.patch("/api/audit/{audit_id}")
async def update_audit_log(audit_id: str, payload: dict):
    allowed_fields = {"raw_input", "intent", "industry", "recommended_tool", "final_prompt", "output"}
    updates = {k: v for k, v in payload.items() if k in allowed_fields}
    if not updates:
        from fastapi import HTTPException
        raise HTTPException(status_code=400, detail="No valid fields to update.")
    conn = get_db()
    set_clause = ", ".join(f"{k} = ?" for k in updates)
    values     = list(updates.values()) + [audit_id]
    conn.execute(f"UPDATE audit_log SET {set_clause} WHERE id = ?", values)
    conn.commit()
    conn.close()
    return {"status": "ok", "audit_id": audit_id}


@router.get("/api/analytics")
async def get_analytics():
    conn = get_db()
    total          = conn.execute("SELECT COUNT(*) as c FROM audit_log").fetchone()["c"]
    intents        = conn.execute("SELECT intent, COUNT(*) as c FROM audit_log GROUP BY intent ORDER BY c DESC").fetchall()
    tools          = conn.execute("SELECT recommended_tool, COUNT(*) as c FROM audit_log GROUP BY recommended_tool ORDER BY c DESC").fetchall()
    industries     = conn.execute("SELECT industry, COUNT(*) as c FROM audit_log GROUP BY industry ORDER BY c DESC LIMIT 5").fetchall()
    avg_rating     = conn.execute("SELECT AVG(rating) as r FROM feedback").fetchone()["r"]
    feedback_count = conn.execute("SELECT COUNT(*) as c FROM feedback").fetchone()["c"]
    issue_types    = conn.execute(
        "SELECT issue_type, COUNT(*) as c FROM feedback WHERE issue_type != '' GROUP BY issue_type ORDER BY c DESC"
    ).fetchall()
    low_rated = conn.execute("""
        SELECT a.intent, a.recommended_tool, f.issue_type, f.comment
        FROM feedback f JOIN audit_log a ON f.audit_id = a.id
        WHERE f.rating <= 2
        ORDER BY f.created_at DESC LIMIT 10
    """).fetchall()
    token_trend = conn.execute(
        "SELECT created_at, token_estimate FROM audit_log ORDER BY created_at DESC LIMIT 10"
    ).fetchall()
    by_user = conn.execute(
        "SELECT user_email, COUNT(*) as c FROM audit_log WHERE user_email != '' AND user_email IS NOT NULL "
        "GROUP BY user_email ORDER BY c DESC LIMIT 20"
    ).fetchall()
    recent_runs = conn.execute(
        "SELECT id, created_at, user_email, raw_input, recommended_tool, intent, policy_blocked "
        "FROM audit_log ORDER BY created_at DESC LIMIT 20"
    ).fetchall()
    conn.close()

    return {
        "total_runs":     total,
        "avg_rating":     round(avg_rating, 1) if avg_rating else None,
        "feedback_count": feedback_count,
        "intents":        [dict(r) for r in intents],
        "tools":          [dict(r) for r in tools],
        "industries":     [dict(r) for r in industries],
        "issue_types":    [dict(r) for r in issue_types],
        "low_rated_runs": [dict(r) for r in low_rated],
        "token_trend":    [dict(r) for r in token_trend],
        "by_user":        [dict(r) for r in by_user],
        "recent_runs":    [dict(r) for r in recent_runs],
    }


# ══════════════════════════════════════════════════════════════════════════════
# ANALYTICS DASHBOARD  (unauthorized, time + role filtered)
# ══════════════════════════════════════════════════════════════════════════════
def _has_column(conn, table: str, column: str) -> bool:
    cols = [r[1] for r in conn.execute(f"PRAGMA table_info({table})").fetchall()]
    return column in cols


def _fill_timeline(rows, since, now, period):
    """Return a continuous list of {label, count} buckets with zero-fill."""
    from datetime import timedelta
    data = dict(rows)
    result = []
    if period == "day":
        current = since.replace(minute=0, second=0, microsecond=0)
        while current <= now:
            bucket  = current.strftime("%H")
            display = current.strftime("%H:00")
            result.append({"label": display, "count": data.get(bucket, 0)})
            current += timedelta(hours=1)
    else:
        current = since.replace(hour=0, minute=0, second=0, microsecond=0)
        while current <= now:
            bucket  = current.strftime("%Y-%m-%d")
            display = current.strftime("%b %d")
            result.append({"label": display, "count": data.get(bucket, 0)})
            current += timedelta(days=1)
    return result


@router.get("/api/analytics-dashboard")
async def get_analytics_dashboard(period: str = "day", role: str = "all"):
    """
    Unauthorized analytics endpoint.
    period: 'day' | 'week' | 'month'
    role:   'all' | any role string (partial LIKE match)
    """
    from datetime import timedelta

    now = datetime.utcnow()
    if period == "week":
        since      = now - timedelta(weeks=1)
        prev_since = since - timedelta(weeks=1)
        tl_fmt     = "%Y-%m-%d"
    elif period == "month":
        since      = now - timedelta(days=30)
        prev_since = since - timedelta(days=30)
        tl_fmt     = "%Y-%m-%d"
    else:  # day (default)
        since      = now - timedelta(days=1)
        prev_since = since - timedelta(days=1)
        tl_fmt     = "%H"

    since_str      = since.isoformat()
    prev_since_str = prev_since.isoformat()

    conn = get_db()

    # Only filter by role if the dedicated role column exists.
    # Never use raw_input as a role fallback — it contains full task descriptions.
    # Build role filter clause safely using parameterized queries
    role_filter_sql  = ""
    role_filter_args = []
    if role != "all" and role.strip():
        role_filter_sql  = " AND LOWER(role) LIKE ?"
        role_filter_args = [f"%{role.lower()}%"]

    base_args      = [since_str]      + role_filter_args
    prev_base_args = [prev_since_str, since_str] + role_filter_args

    # ── Totals ────────────────────────────────────────────────────────────────
    total = conn.execute(
        f"SELECT COUNT(*) as c FROM audit_log WHERE created_at >= ?{role_filter_sql}",
        base_args
    ).fetchone()["c"]

    prev_total = conn.execute(
        f"SELECT COUNT(*) as c FROM audit_log WHERE created_at >= ? AND created_at < ?{role_filter_sql}",
        prev_base_args
    ).fetchone()["c"]

    change_pct = None
    if prev_total and prev_total > 0:
        change_pct = round((total - prev_total) / prev_total * 100)

    # ── By role ───────────────────────────────────────────────────────────────
    by_role_rows = conn.execute(
        "SELECT role, COUNT(*) as count "
        "FROM audit_log "
        "WHERE created_at >= ? "
        "  AND role IS NOT NULL AND TRIM(role) != '' "
        + role_filter_sql +
        " GROUP BY role ORDER BY count DESC LIMIT 15",
        base_args
    ).fetchall()
    by_role = [{"role": r["role"].strip().title() if r["role"].strip().lower() == "general" else r["role"].strip(), "count": r["count"]} for r in by_role_rows]

    # ── By intent ─────────────────────────────────────────────────────────────
    by_intent_rows = conn.execute(
        f"SELECT intent, COUNT(*) as count FROM audit_log "
        f"WHERE created_at >= ? AND (intent IS NOT NULL AND intent != ''){role_filter_sql} "
        f"GROUP BY intent ORDER BY count DESC LIMIT 10",
        base_args
    ).fetchall()
    total_intent = sum(r["count"] for r in by_intent_rows) or 1
    by_intent = [
        {"label": r["intent"] or "—", "count": r["count"],
         "total_pct": round(r["count"] / total_intent * 100)}
        for r in by_intent_rows
    ]

    # ── By tool ───────────────────────────────────────────────────────────────
    by_tool_rows = conn.execute(
        f"SELECT recommended_tool, COUNT(*) as count FROM audit_log "
        f"WHERE created_at >= ? AND (recommended_tool IS NOT NULL AND recommended_tool != ''){role_filter_sql} "
        f"GROUP BY recommended_tool ORDER BY count DESC LIMIT 10",
        base_args
    ).fetchall()
    total_tool = sum(r["count"] for r in by_tool_rows) or 1
    by_tool = [
        {"label": r["recommended_tool"] or "—", "count": r["count"],
         "total_pct": round(r["count"] / total_tool * 100)}
        for r in by_tool_rows
    ]

    # ── Blocked runs ──────────────────────────────────────────────────────────
    blocked = 0
    blocked = conn.execute(
        f"SELECT COUNT(*) as c FROM audit_log "
        f"WHERE created_at >= ? AND policy_blocked = 1{role_filter_sql}",
        base_args
    ).fetchone()["c"]

    # ── Timeline ─────────────────────────────────────────────────────────────
    tl_rows = conn.execute(
        f"SELECT strftime('{tl_fmt}', created_at) as bucket, COUNT(*) as count "
        f"FROM audit_log WHERE created_at >= ?{role_filter_sql} "
        f"GROUP BY bucket ORDER BY bucket ASC",
        base_args
    ).fetchall()
    timeline = _fill_timeline([(r["bucket"], r["count"]) for r in tl_rows], since, now, period)

    conn.close()

    return {
        "period":       period,
        "role_filter":  role,
        "total_runs":   total,
        "change_pct":   change_pct,
        "by_role":      by_role,
        "by_intent":    by_intent,
        "by_tool":      by_tool,
        "blocked_runs": blocked,
        "timeline":     timeline,
    }


# ══════════════════════════════════════════════════════════════════════════════
# PROMPT VERSIONS
# ══════════════════════════════════════════════════════════════════════════════
@router.get("/api/prompt-versions")
async def get_prompt_versions():
    conn = get_db()
    rows = conn.execute(
        "SELECT id, version, intent, industry, change_note, created_at, created_by FROM prompt_versions ORDER BY created_at DESC"
    ).fetchall()
    conn.close()
    return [dict(r) for r in rows]


@router.get("/api/prompt-versions/{version_id}")
async def get_prompt_version(version_id: str):
    conn = get_db()
    row  = conn.execute("SELECT * FROM prompt_versions WHERE id = ?", (version_id,)).fetchone()
    conn.close()
    if not row:
        raise HTTPException(404, "Version not found")
    return dict(row)


@router.post("/api/prompt-versions")
async def create_prompt_version(req: PromptVersionRequest):
    conn = get_db()
    last = conn.execute("SELECT version FROM prompt_versions ORDER BY created_at DESC LIMIT 1").fetchone()
    if last:
        try:
            major, minor = last["version"].split(".")
            new_version  = f"{major}.{int(minor) + 1}"
        except Exception:
            new_version = "1.1"
    else:
        new_version = "1.0"

    vid = str(uuid.uuid4())
    conn.execute(
        "INSERT INTO prompt_versions VALUES (?,?,?,?,?,?,?,?)",
        (vid, new_version, req.intent, req.industry, req.template,
         req.change_note, datetime.utcnow().isoformat(), "user")
    )
    conn.commit()
    conn.close()
    return {"status": "ok", "id": vid, "version": new_version}


@router.post("/api/tools/extract-files")
async def extract_tool_files(files: List[UploadFile] = File(...)):
    """
    Accept one or more PDF / DOCX / PPTX / TXT files for a tool being registered.
    Extract all text, send to LLM, get back a structured JSON matching the tool schema.
    Returns prefilled fields so the UI can populate the registration form.
    """
    combined_text = ""
    for f in files:
        content = await f.read()
        fname   = (f.filename or "").lower()
        try:
            if fname.endswith(".pdf"):
                combined_text += _extract_pdf_text(content) + "\n\n"
            elif fname.endswith(".docx"):
                combined_text += _extract_docx_text(content) + "\n\n"
            elif fname.endswith(".pptx"):
                try:
                    from pptx import Presentation
                    from io import BytesIO
                    prs = Presentation(BytesIO(content))
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                combined_text += shape.text.strip() + "\n"
                    combined_text += "\n"
                except ImportError:
                    combined_text += content.decode("utf-8", errors="ignore") + "\n\n"
            else:
                combined_text += content.decode("utf-8", errors="ignore") + "\n\n"
        except Exception:
            pass

    if not combined_text.strip():
        return {"extracted": {}}

    system_msg = (
        "You are a tool metadata extractor. Given documentation about an AI tool, "
        "extract structured metadata and return ONLY valid JSON — no markdown, no extra text.\n\n"
        "Return this exact JSON shape (leave a field empty/[] if not found):\n"
        "{\n"
        '  "tool_name": "string",\n'
        '  "description": "string — what the tool does, 1-3 sentences",\n'
        '  "category": "string — e.g. AI Writing, Code, Research, Data Analysis",\n'
        '  "url": "string — access URL if mentioned",\n'
        '  "icon": "single emoji that represents this tool",\n'
        '  "best_for": ["list of use-cases this tool is ideal for"],\n'
        '  "strong_signals": ["keywords that strongly indicate this tool should be used"],\n'
        '  "weak_signals": ["secondary keywords"],\n'
        '  "not_for": ["use-cases to avoid with this tool"],\n'
        '  "roles": ["which roles benefit most, e.g. Developer, Business Analyst"],\n'
        '  "output_type": "string — what the tool produces, e.g. report, code, email",\n'
        '  "is_internal": false\n'
        "}"
    )
    user_msg = (
        f"Here is the tool documentation (extracted from uploaded files):\n\n"
        f"{combined_text[:6000]}\n\n"
        "Extract all the fields you can find. For any field you cannot determine, use an empty string or empty array."
    )

    try:
        raw = call_llm(system_msg, user_msg, max_tokens=800, temperature=0.1)
        raw = raw.replace("```json", "").replace("```", "").strip()
        extracted = json.loads(raw)
        return {"extracted": extracted}
    except Exception:
        return {"extracted": {}}


@router.get("/api/tools/registered/{tool_name_or_id}")
async def get_registered_tool(tool_name_or_id: str):
    """Return a single registered tool by tool_name or id (for pre-filling the edit form)."""
    conn = get_db()
    row = conn.execute(
        "SELECT * FROM registered_tools WHERE tool_name = ? OR id = ?",
        (tool_name_or_id, tool_name_or_id),
    ).fetchone()
    conn.close()
    if not row:
        raise HTTPException(404, "Tool not found")
    return dict(row)


@router.post("/api/tools/register")
async def register_tool(req: RegisterToolRequest):
    """Save a new tool (or update existing) into the registered_tools DB table,
    then immediately merge it into the live AI_TOOLS_REGISTRY.

    IMPORTANT: When updating an existing tool, only fields that are non-empty
    in the request will overwrite the stored value — existing content is preserved.
    Pass explicit_edit=True in the request to allow full overwrite (edit mode).
    """
    import service as svc

    tool_id  = str(uuid.uuid4())
    now      = datetime.utcnow().isoformat()

    conn = get_db()
    existing = conn.execute(
        "SELECT * FROM registered_tools WHERE tool_name = ?", (req.tool_name,)
    ).fetchone()

    if existing:
        # Start from whatever raw_data is already stored, then overlay incoming changes
        try:
            existing_raw = json.loads(dict(existing).get("raw_data", "{}") or "{}")
        except Exception:
            existing_raw = {}
        raw_data = {**existing_raw, **(req.raw_data or {})}
    else:
        raw_data = req.raw_data or {}

    if existing:
        ex = dict(existing)
        tool_id = ex["id"]

        is_edit = getattr(req, "explicit_edit", False)

        def _merge_str(new_val, old_val):
            if is_edit:
                return new_val if new_val is not None else old_val
            return new_val if new_val else old_val

        def _merge_list(new_list, old_json):
            if is_edit:
                return new_list if new_list is not None else json.loads(old_json or "[]")
            lst = new_list or []
            return lst if lst else json.loads(old_json or "[]")

        new_desc    = _merge_str(req.description, ex.get("description", ""))
        new_cat     = _merge_str(req.category,    ex.get("category", ""))
        new_url     = _merge_str(req.url,         ex.get("url", ""))
        new_icon    = _merge_str(req.icon,        ex.get("icon", "🤖"))
        new_otype   = _merge_str(req.output_type, ex.get("output_type", ""))
        new_int     = req.is_internal if is_edit else (req.is_internal or bool(ex.get("is_internal", 0)))
        new_best    = _merge_list(req.best_for,        ex.get("best_for",       "[]"))
        new_sig     = _merge_list(req.strong_signals,  ex.get("strong_signals", "[]"))
        new_weak    = _merge_list(req.weak_signals,    ex.get("weak_signals",   "[]"))
        new_notfor  = _merge_list(req.not_for,         ex.get("not_for",        "[]"))
        new_roles   = _merge_list(req.roles,           ex.get("roles",          "[]"))

        changed = {}
        for field, new_v, old_v in [
            ("description", new_desc,  ex.get("description", "")),
            ("category",    new_cat,   ex.get("category", "")),
            ("url",         new_url,   ex.get("url", "")),
            ("icon",        new_icon,  ex.get("icon", "")),
            ("output_type", new_otype, ex.get("output_type", "")),
        ]:
            if str(new_v) != str(old_v):
                changed[field] = {"from": old_v, "to": new_v}
        for field, new_list, old_json in [
            ("best_for",       new_best,   ex.get("best_for",       "[]")),
            ("strong_signals", new_sig,    ex.get("strong_signals", "[]")),
            ("weak_signals",   new_weak,   ex.get("weak_signals",   "[]")),
            ("not_for",        new_notfor, ex.get("not_for",        "[]")),
            ("roles",          new_roles,  ex.get("roles",          "[]")),
        ]:
            old_list = json.loads(old_json) if isinstance(old_json, str) else (old_json or [])
            if sorted(new_list) != sorted(old_list):
                changed[field] = {"from": old_list, "to": new_list}

        conn.execute(
            """UPDATE registered_tools SET
               description=?, category=?, url=?, icon=?,
               best_for=?, strong_signals=?, weak_signals=?, not_for=?,
               roles=?, output_type=?, is_internal=?, raw_data=?, updated_at=?
               WHERE tool_name=?""",
            (
                new_desc, new_cat, new_url, new_icon,
                json.dumps(new_best), json.dumps(new_sig),
                json.dumps(new_weak), json.dumps(new_notfor),
                json.dumps(new_roles), new_otype,
                1 if new_int else 0, json.dumps(raw_data), now,
                req.tool_name,
            ),
        )
        conn.commit()
        conn.close()
        log_tool_change(req.tool_name, "updated", changed)
    else:
        conn.execute(
            """INSERT INTO registered_tools
               (id, tool_name, description, category, url, icon,
                best_for, strong_signals, weak_signals, not_for,
                roles, output_type, is_internal, raw_data, created_at, updated_at)
               VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)""",
            (
                tool_id, req.tool_name,
                req.description, req.category, req.url, req.icon,
                json.dumps(req.best_for or []),
                json.dumps(req.strong_signals or []),
                json.dumps(req.weak_signals or []),
                json.dumps(req.not_for or []),
                json.dumps(req.roles or []),
                req.output_type, 1 if req.is_internal else 0,
                json.dumps(raw_data), now, now,
            ),
        )
        conn.commit()
        conn.close()
        log_tool_change(req.tool_name, "registered", {
            "description": req.description,
            "category": req.category,
        })

    svc._merge_db_tools_into_registry()

    return {"status": "ok", "tool_id": tool_id, "tool_name": req.tool_name}


@router.get("/api/tools/registered")
async def list_registered_tools():
    """Return all tools stored in the registered_tools DB table."""
    conn = get_db()
    rows = conn.execute(
        "SELECT * FROM registered_tools ORDER BY created_at DESC"
    ).fetchall()
    conn.close()
    return [dict(r) for r in rows]


@router.delete("/api/tools/registered/{tool_id}")
async def delete_registered_tool(tool_id: str):
    """Delete a DB-registered tool and remove it from the live registry."""
    import service as svc
    conn = get_db()
    row = conn.execute(
        "SELECT tool_name FROM registered_tools WHERE id = ?", (tool_id,)
    ).fetchone()
    if not row:
        raise HTTPException(404, "Tool not found")
    tool_name = row["tool_name"]
    conn.execute("DELETE FROM registered_tools WHERE id = ?", (tool_id,))
    conn.commit()
    conn.close()
    svc.AI_TOOLS_REGISTRY.pop(tool_name, None)
    log_tool_change(tool_name, "deleted")
    return {"status": "ok"}


@router.get("/api/tool-change-log")
async def get_tool_change_log(page: int = 1, per_page: int = 30, tool_name: str = ""):
    """Return paginated tool change log entries.
    Excludes file_deleted entries and updated entries with no changed fields
    so each user action appears as exactly one log row.
    """
    conn = get_db()
    filters = [
        "action != 'file_deleted'",
        "NOT (action = 'updated' AND (changed_fields IS NULL OR changed_fields = '{}' OR changed_fields = ''))",
    ]
    args = []
    if tool_name.strip():
        filters.append("tool_name LIKE ?")
        args.append(f"%{tool_name.strip()}%")
    where = "WHERE " + " AND ".join(filters)
    base = f"FROM tool_change_log {where}"
    total = conn.execute(f"SELECT COUNT(*) as c {base}", args).fetchone()["c"]
    offset = (page - 1) * per_page
    rows = conn.execute(
        f"SELECT * {base} ORDER BY created_at DESC LIMIT ? OFFSET ?",
        args + [per_page, offset],
    ).fetchall()
    conn.close()
    return {
        "total":    total,
        "page":     page,
        "per_page": per_page,
        "logs":     [dict(r) for r in rows],
    }


@router.delete("/api/tool-docs/file")
async def delete_tool_doc_file(filename: str, tool_name: str = ""):
    """Delete all ChromaDB chunks that came from a specific source file."""
    try:
        results = tool_knowledge_collection.get(
            where={"source_file": filename}, include=[]
        )
        ids = results.get("ids", [])
        if ids:
            tool_knowledge_collection.delete(ids=ids)
        log_tool_change(
            tool_name=tool_name or filename,
            action="updated",
            changed_fields={"file_removed": filename, "chunks_removed": len(ids)},
        )
        return {"status": "ok", "deleted_chunks": len(ids), "filename": filename}
    except Exception as e:
        raise HTTPException(500, str(e))


@router.post("/api/tool-docs/upload/{tool_name}")
async def upload_tool_docs_for_tool(tool_name: str, files: List[UploadFile] = File(...)):
    """
    Upload documents for a specific tool — every chunk is tagged directly with
    that tool_name, bypassing the LLM classifier entirely.
    """
    import service as svc

    processed_files = []
    total_chunks    = 0

    for f in files:
        content  = await f.read()
        filename = f.filename or "unknown"
        try:
            summary = svc.ingest_tool_document_direct(content, filename, tool_name)
            count   = summary.get(tool_name, 0)
            total_chunks += count
            processed_files.append({"filename": filename, "chunks": count})
        except Exception as exc:
            processed_files.append({"filename": filename, "error": str(exc)})

    filenames = [f["filename"] for f in processed_files if "error" not in f]
    if filenames:
        log_tool_change(
            tool_name=tool_name,
            action="updated",
            changed_fields={
                "files_uploaded": ", ".join(filenames),
                "chunks_added":   total_chunks,
            },
        )

    return {
        "status":          "ok",
        "tool_name":       tool_name,
        "files_processed": len(processed_files),
        "files":           processed_files,
        "total_chunks":    total_chunks,
    }


@router.post("/api/tool-docs/upload")
async def upload_tool_docs(files: List[UploadFile] = File(...)):
    """
    Upload one or more PDF/DOCX/PPTX/TXT documents.
    Each file is chunked, each chunk is auto-classified to a tool via LLM,
    and all chunks are stored in the tool_knowledge ChromaDB collection.
    No manual labelling required — the agent figures out which tool each chunk belongs to.
    """
    import service as svc

    tool_names = list(svc.AI_TOOLS_REGISTRY.keys())
    if not tool_names:
        raise HTTPException(400, "No tools in registry. Upload your Excel first so the classifier knows what tools exist.")

    overall_summary = {}
    processed_files = []

    for f in files:
        content  = await f.read()
        filename = f.filename or "unknown"
        try:
            file_summary = svc.ingest_tool_document(content, filename, tool_names)
            for tool, count in file_summary.items():
                overall_summary[tool] = overall_summary.get(tool, 0) + count
            processed_files.append({
                "filename":    filename,
                "chunks":      sum(file_summary.values()),
                "classified":  {k: v for k, v in file_summary.items() if k != "unclassified"},
                "unclassified": file_summary.get("unclassified", 0),
            })
        except Exception as exc:
            processed_files.append({"filename": filename, "error": str(exc)})

    return {
        "status":          "ok",
        "files_processed": len(processed_files),
        "files":           processed_files,
        "tool_summary":    {k: v for k, v in overall_summary.items() if k != "unclassified"},
        "unclassified_chunks": overall_summary.get("unclassified", 0),
    }


@router.get("/api/tool-docs/status")
async def tool_docs_status():
    """Return a summary of all tool knowledge chunks stored in ChromaDB."""
    return {"status": get_tool_knowledge_status()}


@router.delete("/api/tool-docs/clear")
async def clear_tool_docs():
    """Delete all tool knowledge chunks from ChromaDB (admin reset)."""
    try:
        all_ids = tool_knowledge_collection.get()["ids"]
        if all_ids:
            tool_knowledge_collection.delete(ids=all_ids)
        return {"status": "ok", "deleted": len(all_ids)}
    except Exception as e:
        raise HTTPException(500, str(e))


@router.delete("/api/tool-docs/clear/{tool_name}")
async def clear_tool_docs_for_tool(tool_name: str):
    """Delete all knowledge chunks for a specific tool."""
    try:
        results = tool_knowledge_collection.get(
            where={"tool_name": tool_name}, include=[]
        )
        ids = results.get("ids", [])
        if ids:
            tool_knowledge_collection.delete(ids=ids)
        return {"status": "ok", "deleted": len(ids), "tool_name": tool_name}
    except Exception as e:
        raise HTTPException(500, str(e))


@router.post("/api/upload-tools-registry")
async def upload_tools_registry(file: UploadFile = File(...)):
    import service  # import the module, not the dict

    filename = file.filename or ""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext not in ("xlsx", "xlsm", "xls"):
        raise HTTPException(400, "Only Excel files (.xlsx, .xlsm, .xls) are supported")

    content = await file.read()

    try:
        service.reload_tools_registry(excel_bytes=content)
    except ValueError as e:
        raise HTTPException(400, str(e))
    except Exception as e:
        raise HTTPException(400, f"Could not read Excel file: {str(e)}")

    if not service.AI_TOOLS_REGISTRY:
        raise HTTPException(400, "File was read but no tools were found. Ensure your sheet has a tool name column and at least one data row.")

    return {"status": "ok", "tools_loaded": len(service.AI_TOOLS_REGISTRY)}


# ══════════════════════════════════════════════════════════════════════════════
# SCENARIO LIBRARY ENDPOINTS
# ══════════════════════════════════════════════════════════════════════════════

@router.get("/api/scenarios")
async def get_scenarios():
    """Return the full scenario library as a JSON list."""
    import service as _svc
    return {"status": "ok", "scenarios": _svc.SCENARIO_LIBRARY, "count": len(_svc.SCENARIO_LIBRARY)}


@router.post("/api/scenarios/register")
async def register_scenario(req: RegisterScenarioRequest):
    """Add a single scenario to the in-memory SCENARIO_LIBRARY."""
    import service as _svc
    new_scenario = {
        "mega_group":     req.mega_group.strip(),
        "category":       req.category.strip() if req.category else "",
        "title":          req.title.strip(),
        "persona":        req.persona.strip() if req.persona else "",
        "activate_phase": req.activate_phase.strip() if req.activate_phase else "",
        "scenario":       req.scenario.strip(),
    }
    _svc.SCENARIO_LIBRARY.append(new_scenario)
    return {"status": "ok", "scenarios_loaded": len(_svc.SCENARIO_LIBRARY)}


# ══════════════════════════════════════════════════════════════════════════════
# SCENARIO SUGGESTIONS
# ══════════════════════════════════════════════════════════════════════════════

@router.post("/api/scenario-suggestions/submit")
async def submit_scenario_suggestion(req: RegisterScenarioRequest, submitted_by: str = ""):
    """Submit a scenario suggestion for admin review."""
    from service import get_db
    suggestion_id = str(uuid.uuid4())
    submitted_at  = datetime.utcnow().isoformat()
    conn = get_db()
    conn.execute(
        """INSERT INTO scenario_suggestions
           (id, title, mega_group, category, persona, activate_phase, scenario, submitted_by, submitted_at, status, admin_note, reviewed_at)
           VALUES (?,?,?,?,?,?,?,?,?,'pending','','')""",
        (
            suggestion_id,
            req.title.strip(),
            req.mega_group.strip(),
            req.category.strip() if req.category else "",
            req.persona.strip() if req.persona else "",
            req.activate_phase.strip() if req.activate_phase else "",
            req.scenario.strip(),
            submitted_by.strip(),
            submitted_at,
        ),
    )
    conn.commit()
    conn.close()
    return {"status": "ok", "suggestion_id": suggestion_id}


@router.get("/api/scenario-suggestions")
async def list_scenario_suggestions(
    status: str = "all",
    search: str = "",
    page: int = 1,
    per_page: int = 20,
):
    """List scenario suggestions with optional status filter and search."""
    from service import get_db
    conn = get_db()

    where_parts = []
    params: list = []

    if status != "all":
        where_parts.append("status = ?")
        params.append(status)

    if search.strip():
        q = f"%{search.strip().lower()}%"
        where_parts.append("(LOWER(title) LIKE ? OR LOWER(mega_group) LIKE ? OR LOWER(scenario) LIKE ? OR LOWER(submitted_by) LIKE ?)")
        params.extend([q, q, q, q])

    where_sql = ("WHERE " + " AND ".join(where_parts)) if where_parts else ""

    total = conn.execute(f"SELECT COUNT(*) as c FROM scenario_suggestions {where_sql}", params).fetchone()["c"]

    offset = (page - 1) * per_page
    rows = conn.execute(
        f"SELECT * FROM scenario_suggestions {where_sql} ORDER BY submitted_at DESC LIMIT ? OFFSET ?",
        params + [per_page, offset],
    ).fetchall()
    conn.close()

    return {
        "total":    total,
        "page":     page,
        "per_page": per_page,
        "items":    [dict(r) for r in rows],
    }


@router.post("/api/scenario-suggestions/{suggestion_id}/approve")
async def approve_scenario_suggestion(suggestion_id: str, admin_note: str = ""):
    """Approve a suggestion and push it into the live SCENARIO_LIBRARY."""
    from service import get_db, SCENARIO_LIBRARY
    conn = get_db()
    row = conn.execute("SELECT * FROM scenario_suggestions WHERE id = ?", (suggestion_id,)).fetchone()
    if not row:
        raise HTTPException(404, "Suggestion not found")

    conn.execute(
        "UPDATE scenario_suggestions SET status='approved', admin_note=?, reviewed_at=? WHERE id=?",
        (admin_note, datetime.utcnow().isoformat(), suggestion_id),
    )
    conn.commit()
    conn.close()

    SCENARIO_LIBRARY.append({
        "mega_group":     row["mega_group"],
        "category":       row["category"],
        "title":          row["title"],
        "persona":        row["persona"],
        "activate_phase": row["activate_phase"],
        "scenario":       row["scenario"],
    })
    return {"status": "ok", "scenarios_loaded": len(SCENARIO_LIBRARY)}


@router.post("/api/scenario-suggestions/{suggestion_id}/reject")
async def reject_scenario_suggestion(suggestion_id: str, admin_note: str = ""):
    """Reject a suggestion."""
    from service import get_db
    conn = get_db()
    row = conn.execute("SELECT id FROM scenario_suggestions WHERE id = ?", (suggestion_id,)).fetchone()
    if not row:
        raise HTTPException(404, "Suggestion not found")
    conn.execute(
        "UPDATE scenario_suggestions SET status='rejected', admin_note=?, reviewed_at=? WHERE id=?",
        (admin_note, datetime.utcnow().isoformat(), suggestion_id),
    )
    conn.commit()
    conn.close()
    return {"status": "ok"}


@router.post("/api/upload-scenario-library")
async def upload_scenario_library(file: UploadFile = File(...)):
    """Upload a new Scenario Library Excel file and reload SCENARIO_LIBRARY in-place."""
    import service as _svc

    filename = file.filename or ""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext not in ("xlsx", "xlsm", "xls"):
        raise HTTPException(400, "Only Excel files (.xlsx, .xlsm, .xls) are supported")

    content = await file.read()

    try:
        _svc.reload_scenario_library(excel_bytes=content)
    except ValueError as e:
        raise HTTPException(400, str(e))
    except Exception as e:
        raise HTTPException(400, f"Could not read Excel file: {str(e)}")

    if not _svc.SCENARIO_LIBRARY:
        raise HTTPException(400, "File was read but no scenarios were found. Ensure your sheet has a title/scenario column and at least one data row.")

    return {"status": "ok", "scenarios_loaded": len(_svc.SCENARIO_LIBRARY)}


@router.patch("/api/audit/{audit_id}")
async def update_audit_log(audit_id: str, req: AuditUpdateRequest):
    conn = get_db()
    row = conn.execute("SELECT * FROM audit_log WHERE id = ?", (audit_id,)).fetchone()
    if not row:
        conn.close()
        raise HTTPException(404, "Audit log not found")

    current = dict(row)

    new_raw_input    = req.raw_input    if req.raw_input is not None else current.get("raw_input", "")
    new_final_prompt = req.final_prompt if req.final_prompt is not None else current.get("final_prompt", "")
    new_output       = req.output       if req.output is not None else current.get("output", "")

    conn.execute(
        "UPDATE audit_log SET raw_input = ?, final_prompt = ?, output = ? WHERE id = ?",
        (new_raw_input, new_final_prompt, new_output, audit_id),
    )
    conn.commit()

    updated = conn.execute("SELECT * FROM audit_log WHERE id = ?", (audit_id,)).fetchone()
    conn.close()
    return dict(updated)



# ══════════════════════════════════════════════════════════════════════════════
# CLARIFIER AGENT  —  Agent 1: conversational intent gathering
#
# Uses the same Clarifier Agent prompt from Agentic AI Navigator.
# Collects exactly 3 pieces of information one question at a time:
#   1. Role  2. Core Task  3. Specific Parameters
# Emits [SATISFIED] when all 3 are clear, which triggers action='ready'.
# ══════════════════════════════════════════════════════════════════════════════

_CLARIFIER_SYSTEM_PROMPT = """
You are an expert AI Requirements Gatherer. Your job is to help users clarify their task so you can recommend the right AI tools later.
You MUST ensure you have three pieces of information from the user:
1. Their Role (e.g., Developer, Marketer, HR, Student)
2. Their Core Task (e.g., Write code, summarize documents, analyze data)
3. Specific Parameters (e.g., How many documents? What programming language? What expected output format?)

INSTRUCTIONS:
- If ANY of these 3 elements are missing or vague, politely ask ONE clear follow-up question to get the missing info. Do not ask multiple questions at once.
- Be professional, conversational, and helpful.
- Keep responses SHORT and SCANNABLE — no long paragraphs.
- If the user says 'skip', 'proceed', 'generate', or similar — treat the available info as sufficient and output the [SATISFIED] block immediately.
- If ALL 3 elements are clearly provided, you are satisfied. You must then output ONLY the following format:
[SATISFIED]
Role: <user_role>
Task Details: <detailed_task_description_with_parameters>

Do not add any conversational text if you are satisfied. Just output the [SATISFIED] block.
"""


def _parse_satisfied_block(text: str, fallback_role: str, fallback_task: str) -> dict:
    """
    Parse a [SATISFIED] block from the Clarifier Agent response.
    Returns { role, task_type, task_description } ready for downstream nodes.
    """
    lines = text.replace("[SATISFIED]", "").strip().splitlines()
    role_val = fallback_role or "general"
    task_val = fallback_task or ""

    for line in lines:
        if line.lower().startswith("role:"):
            role_val = line.split(":", 1)[1].strip() or role_val
        elif line.lower().startswith("task details:"):
            task_val = line.split(":", 1)[1].strip() or task_val

    TASK_TYPE_KEYWORDS = {
        "Research & Analysis":  ["research", "analys", "findings", "review", "report"],
        "Writing & Docs":       ["write", "document", "draft", "proposal", "summary"],
        "Strategy & Planning":  ["strategy", "plan", "roadmap", "decision"],
        "Data Analysis":        ["data", "dashboard", "kpi", "metric", "chart", "insight"],
        "Code & Dev":           ["code", "script", "debug", "develop", "program", "api", "automate"],
        "Creative Content":     ["blog", "article", "creative", "post", "copy", "marketing"],
        "Communication":        ["email", "message", "communicate", "reply"],
        "Learning & Training":  ["learn", "training", "tutorial", "course"],
        "Process Automation":   ["automate", "workflow", "process", "pipeline"],
        "Decision Support":     ["decide", "compare", "evaluate", "recommend"],
    }
    detected_task_type = "general"
    task_lower = task_val.lower()
    for tt, kws in TASK_TYPE_KEYWORDS.items():
        if any(kw in task_lower for kw in kws):
            detected_task_type = tt
            break

    return {
        "role":             role_val,
        "task_type":        detected_task_type,
        "task_description": task_val,
    }


@router.post("/api/chat-gather")
async def chat_gather(req: ChatGatherRequest):
    """
    Clarifier Agent (Agent 1) — conversational intent gathering.
    Builds a proper OpenAI message array (system + full conversation history)
    so the LLM sees every turn and responds contextually, not with a fixed reply.
    """
    messages = [{"role": "system", "content": _CLARIFIER_SYSTEM_PROMPT}]

    # Normalise frontend role 'agent' → 'assistant' so the LLM understands turns
    def _normalise_role(r: str) -> str:
        return "assistant" if r in ("agent", "assistant", "bot") else "user"

    # If there is no conversation history yet, the user_input is the opening message
    if req.messages:
        for m in req.messages:
            messages.append({"role": _normalise_role(m.role), "content": m.content})
    elif req.user_input:
        messages.append({"role": "user", "content": req.user_input})

    try:
        raw = call_llm_messages(messages, max_tokens=400, temperature=0.2)
        raw = (raw or "").strip()

        if not raw:
            raise ValueError("Empty LLM response")

        if "[SATISFIED]" in raw:
            parsed = _parse_satisfied_block(raw, req.role or "general", req.user_input or "")
            return {
                "action":           "ready",
                "role":             parsed["role"],
                "task_type":        parsed["task_type"],
                "task_description": parsed["task_description"],
                "message":          None,
            }

        return {
            "action":           "question",
            "message":          raw,
            "role":             None,
            "task_type":        None,
            "task_description": None,
        }
    except Exception as e:
        import logging
        logging.getLogger("routes").error("Clarifier Agent error: %s", e, exc_info=True)
        return {
            "action":           "question",
            "message":          "I'm here to help! Could you describe what you'd like to accomplish today?",
            "role":             None,
            "task_type":        None,
            "task_description": None,
        }


@router.post("/api/chat-summarize")
async def chat_summarize(req: ChatSummarizeRequest):
    """
    Force-summarize the conversation using the Clarifier Agent even if [SATISFIED]
    was not emitted (used for 'Skip All'). Sends the full message array with a
    final user instruction to emit [SATISFIED] immediately with what is known.
    """
    messages = [{"role": "system", "content": _CLARIFIER_SYSTEM_PROMPT}]

    def _normalise_role(r: str) -> str:
        return "assistant" if r in ("agent", "assistant", "bot") else "user"

    if req.user_input:
        messages.append({"role": "user", "content": req.user_input})

    for m in req.messages:
        messages.append({"role": _normalise_role(m.role), "content": m.content})

    messages.append({
        "role": "user",
        "content": (
            "I'd like to skip further questions. "
            "Please output the [SATISFIED] block now using the best available information."
        ),
    })

    try:
        raw = call_llm_messages(messages, max_tokens=300, temperature=0.1)
        raw = raw.strip()

        if "[SATISFIED]" in raw:
            parsed = _parse_satisfied_block(raw, req.role or "general", req.user_input or "")
            return parsed

        raw_json = call_llm(
            "You are a task summarizer. Return ONLY valid JSON — no markdown, no extra text.",
            (
                (f"Initial description: \"{req.user_input}\"\n\n" if req.user_input else "") +
                "Conversation:\n" +
                "\n".join(f"{'User' if m['role'] == 'user' else 'Agent'}: {m['content']}" for m in messages[1:]) +
                "\n\nExtract and return JSON:\n"
                "{\"role\": \"...\", \"task_type\": \"...\", \"task_description\": \"...\"}"
            ),
            max_tokens=300,
            temperature=0.1,
        )
        raw_json = raw_json.replace("```json", "").replace("```", "").strip()
        data = json.loads(raw_json)
        return {
            "role":             data.get("role", req.role or "general"),
            "task_type":        data.get("task_type", req.task_type or "general"),
            "task_description": data.get("task_description", req.user_input or ""),
        }
    except Exception:
        return {
            "role":             req.role or "general",
            "task_type":        req.task_type or "general",
            "task_description": req.user_input or "",
        }


from pathlib import Path
from fastapi.responses import Response

@router.get("/saml/metadata")
async def saml_metadata():
    cert = Path("saml/sp.crt").read_text()
    cert = cert.replace("-----BEGIN CERTIFICATE-----", "") \
               .replace("-----END CERTIFICATE-----", "") \
               .replace("\n", "").strip()

    base_url = "https://ai-navigator-ashpbzhbcmgeerbt.northeurope-01.azurewebsites.net/"  # ← local for now

    xml = f"""<?xml version="1.0"?>
<md:EntityDescriptor
    xmlns:md="urn:oasis:names:tc:SAML:2.0:metadata"
    xmlns:ds="http://www.w3.org/2000/09/xmldsig#"
    entityID="{base_url}/saml/metadata">
  <md:SPSSODescriptor
      AuthnRequestsSigned="false"
      WantAssertionsSigned="true"
      protocolSupportEnumeration="urn:oasis:names:tc:SAML:2.0:protocol">
    <md:KeyDescriptor use="signing">
      <ds:KeyInfo>
        <ds:X509Data>
          <ds:X509Certificate>{cert}</ds:X509Certificate>
        </ds:X509Data>
      </ds:KeyInfo>
    </md:KeyDescriptor>
    <md:NameIDFormat>urn:oasis:names:tc:SAML:1.1:nameid-format:emailAddress</md:NameIDFormat>
    <md:AssertionConsumerService
        Binding="urn:oasis:names:tc:SAML:2.0:bindings:HTTP-POST"
        Location="{base_url}/saml/acs"
        index="1"/>
  </md:SPSSODescriptor>
</md:EntityDescriptor>"""
    return Response(content=xml, media_type="application/xml")


# ══════════════════════════════════════════════════════════════════════════════
# AUTH ROUTES  — identity, permissions, admin management
# ══════════════════════════════════════════════════════════════════════════════

@router.post("/api/auth/identify")
async def identify(email: str = Form(...)):
    """
    Called when the user enters their email on the login screen.
    Returns their role and the list of permitted features.
    """
    try:
        result = _auth.identify_user(email)
        result["permissions"] = _auth.get_permissions(result["role"])
        return result
    except ValueError as e:
        raise HTTPException(400, str(e))
    except Exception as e:
        raise HTTPException(500, f"Auth error: {e}")


@router.get("/api/auth/admins")
async def get_admins():
    """List all admins in NavigatorAdmins."""
    try:
        return {"admins": _auth.list_admins()}
    except Exception as e:
        raise HTTPException(500, str(e))


@router.post("/api/auth/admins/add")
async def add_admin(email: str = Form(...), name: str = Form("")):
    """Add an email to NavigatorAdmins."""
    try:
        return _auth.add_admin(email, name)
    except Exception as e:
        raise HTTPException(500, str(e))


@router.post("/api/auth/admins/remove")
async def remove_admin(email: str = Form(...)):
    """Remove an email from NavigatorAdmins."""
    try:
        return _auth.remove_admin(email)
    except Exception as e:
        raise HTTPException(500, str(e))


@router.get("/api/auth/users")
async def get_users(page: int = 1, per_page: int = 50):
    """List all users in NavigatorUsers."""
    try:
        return _auth.list_users(page, per_page)
    except Exception as e:
        raise HTTPException(500, str(e))