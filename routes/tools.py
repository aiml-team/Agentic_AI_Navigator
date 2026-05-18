import json
import uuid
from datetime import datetime
from typing import List

from fastapi import APIRouter, HTTPException, UploadFile, File
from schemas import RegisterToolRequest
from services.registry import AI_TOOLS_REGISTRY, _merge_db_tools_into_registry, reload_tools_registry
from services.database import get_db, log_tool_change
from services.chromadb_store import (
    tool_knowledge_collection,
    ingest_tool_document,
    ingest_tool_document_direct,
    get_tool_knowledge_status,
)
from services.llm_client import call_llm
from services.file_utils import extract_pdf_text as _extract_pdf_text, extract_docx_text as _extract_docx_text

router = APIRouter()


@router.get("/api/tools")
async def get_tools():
    return {
        name: {
            "description": info.get("description") or "",
            "desc_content": (info.get("raw_data") or {}).get("Desc_Content", ""),
            "category":    info.get("category", ""),
            "icon":        info.get("icon", ""),
            "best_for":    info.get("best_for") or [],
            "not_for":     info.get("not_for") or [],
            "strong_signals": info.get("strong_signals") or [],
            "weak_signals":   info.get("weak_signals") or [],
            "roles":       info.get("roles") or [],
            "output_type": info.get("output_type", ""),
            "url":         info.get("url", ""),
            "is_internal": info.get("is_internal", False),
            "raw_data":    info.get("raw_data", {}),
            "_source":     info.get("_source", "registry"),
        }
        for name, info in AI_TOOLS_REGISTRY.items()
    }


@router.post("/api/tools/extract-files")
async def extract_tool_files(files: List[UploadFile] = File(...)):
    combined_text = ""
    for f in files:
        content = await f.read()
        fname   = (f.filename or "").lower()
        try:
            if fname.endswith(".pdf"):
                combined_text += _extract_pdf_text(content) + "\n\n"
            elif fname.endswith(".docx"):
                combined_text += _extract_docx_text(content) + "\n\n"
            elif fname.endswith(".pptx"):
                try:
                    from pptx import Presentation
                    from io import BytesIO
                    prs = Presentation(BytesIO(content))
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                combined_text += shape.text.strip() + "\n"
                    combined_text += "\n"
                except ImportError:
                    combined_text += content.decode("utf-8", errors="ignore") + "\n\n"
            else:
                combined_text += content.decode("utf-8", errors="ignore") + "\n\n"
        except Exception:
            pass

    if not combined_text.strip():
        return {"extracted": {}}

    system_msg = (
        "You are a tool metadata extractor. Given documentation about an AI tool, "
        "extract structured metadata and return ONLY valid JSON — no markdown, no extra text.\n\n"
        "Return this exact JSON shape (leave a field empty/[] if not found):\n"
        "{\n"
        '  "tool_name": "string",\n'
        '  "description": "string — what the tool does, 1-3 sentences",\n'
        '  "category": "string — e.g. AI Writing, Code, Research, Data Analysis",\n'
        '  "url": "string — access URL if mentioned",\n'
        '  "icon": "single emoji that represents this tool",\n'
        '  "best_for": ["list of use-cases this tool is ideal for"],\n'
        '  "strong_signals": ["keywords that strongly indicate this tool should be used"],\n'
        '  "weak_signals": ["secondary keywords"],\n'
        '  "not_for": ["use-cases to avoid with this tool"],\n'
        '  "roles": ["which roles benefit most, e.g. Developer, Business Analyst"],\n'
        '  "output_type": "string — what the tool produces, e.g. report, code, email",\n'
        '  "is_internal": false\n'
        "}"
    )
    user_msg = (
        f"Here is the tool documentation (extracted from uploaded files):\n\n"
        f"{combined_text[:6000]}\n\n"
        "Extract all the fields you can find. For any field you cannot determine, use an empty string or empty array."
    )

    try:
        raw = call_llm(system_msg, user_msg, max_tokens=800, temperature=0.1)
        raw = raw.replace("```json", "").replace("```", "").strip()
        extracted = json.loads(raw)
        return {"extracted": extracted}
    except Exception:
        return {"extracted": {}}


@router.get("/api/tools/registered/{tool_name_or_id}")
async def get_registered_tool(tool_name_or_id: str):
    conn = get_db()
    row = conn.execute(
        "SELECT * FROM registered_tools WHERE tool_name = ? OR id = ?",
        (tool_name_or_id, tool_name_or_id),
    ).fetchone()
    conn.close()
    if not row:
        raise HTTPException(404, "Tool not found")
    return dict(row)


@router.post("/api/tools/register")
async def register_tool(req: RegisterToolRequest):
    tool_id  = str(uuid.uuid4())
    now      = datetime.utcnow().isoformat()

    conn = get_db()
    existing = conn.execute(
        "SELECT * FROM registered_tools WHERE tool_name = ?", (req.tool_name,)
    ).fetchone()

    if existing:
        try:
            existing_raw = json.loads(dict(existing).get("raw_data", "{}") or "{}")
        except Exception:
            existing_raw = {}
        raw_data = {**existing_raw, **(req.raw_data or {})}
    else:
        raw_data = req.raw_data or {}

    if existing:
        ex = dict(existing)
        tool_id = ex["id"]

        is_edit = getattr(req, "explicit_edit", False)

        def _merge_str(new_val, old_val):
            if is_edit:
                return new_val if new_val is not None else old_val
            return new_val if new_val else old_val

        def _merge_list(new_list, old_json):
            if is_edit:
                return new_list if new_list is not None else json.loads(old_json or "[]")
            lst = new_list or []
            return lst if lst else json.loads(old_json or "[]")

        new_desc    = _merge_str(req.description, ex.get("description", ""))
        new_cat     = _merge_str(req.category,    ex.get("category", ""))
        new_url     = _merge_str(req.url,         ex.get("url", ""))
        new_icon    = _merge_str(req.icon,        ex.get("icon", "🤖"))
        new_otype   = _merge_str(req.output_type, ex.get("output_type", ""))
        new_int     = req.is_internal if is_edit else (req.is_internal or bool(ex.get("is_internal", 0)))
        new_best    = _merge_list(req.best_for,        ex.get("best_for",       "[]"))
        new_sig     = _merge_list(req.strong_signals,  ex.get("strong_signals", "[]"))
        new_weak    = _merge_list(req.weak_signals,    ex.get("weak_signals",   "[]"))
        new_notfor  = _merge_list(req.not_for,         ex.get("not_for",        "[]"))
        new_roles   = _merge_list(req.roles,           ex.get("roles",          "[]"))

        changed = {}
        for field, new_v, old_v in [
            ("description", new_desc,  ex.get("description", "")),
            ("category",    new_cat,   ex.get("category", "")),
            ("url",         new_url,   ex.get("url", "")),
            ("icon",        new_icon,  ex.get("icon", "")),
            ("output_type", new_otype, ex.get("output_type", "")),
        ]:
            if str(new_v) != str(old_v):
                changed[field] = {"from": old_v, "to": new_v}
        for field, new_list, old_json in [
            ("best_for",       new_best,   ex.get("best_for",       "[]")),
            ("strong_signals", new_sig,    ex.get("strong_signals", "[]")),
            ("weak_signals",   new_weak,   ex.get("weak_signals",   "[]")),
            ("not_for",        new_notfor, ex.get("not_for",        "[]")),
            ("roles",          new_roles,  ex.get("roles",          "[]")),
        ]:
            old_list = json.loads(old_json) if isinstance(old_json, str) else (old_json or [])
            if sorted(new_list) != sorted(old_list):
                changed[field] = {"from": old_list, "to": new_list}

        conn.execute(
            """UPDATE registered_tools SET
               description=?, category=?, url=?, icon=?,
               best_for=?, strong_signals=?, weak_signals=?, not_for=?,
               roles=?, output_type=?, is_internal=?, raw_data=?, updated_at=?
               WHERE tool_name=?""",
            (
                new_desc, new_cat, new_url, new_icon,
                json.dumps(new_best), json.dumps(new_sig),
                json.dumps(new_weak), json.dumps(new_notfor),
                json.dumps(new_roles), new_otype,
                1 if new_int else 0, json.dumps(raw_data), now,
                req.tool_name,
            ),
        )
        conn.commit()
        conn.close()
        log_tool_change(req.tool_name, "updated", changed)
    else:
        conn.execute(
            """INSERT INTO registered_tools
               (id, tool_name, description, category, url, icon,
                best_for, strong_signals, weak_signals, not_for,
                roles, output_type, is_internal, raw_data, created_at, updated_at)
               VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?)""",
            (
                tool_id, req.tool_name,
                req.description, req.category, req.url, req.icon,
                json.dumps(req.best_for or []),
                json.dumps(req.strong_signals or []),
                json.dumps(req.weak_signals or []),
                json.dumps(req.not_for or []),
                json.dumps(req.roles or []),
                req.output_type, 1 if req.is_internal else 0,
                json.dumps(raw_data), now, now,
            ),
        )
        conn.commit()
        conn.close()
        log_tool_change(req.tool_name, "registered", {
            "description": req.description,
            "category": req.category,
        })

    _merge_db_tools_into_registry()

    return {"status": "ok", "tool_id": tool_id, "tool_name": req.tool_name}


@router.get("/api/tools/registered")
async def list_registered_tools():
    conn = get_db()
    rows = conn.execute(
        "SELECT * FROM registered_tools ORDER BY created_at DESC"
    ).fetchall()
    conn.close()
    return [dict(r) for r in rows]


@router.delete("/api/tools/registered/{tool_id}")
async def delete_registered_tool(tool_id: str):
    conn = get_db()
    row = conn.execute(
        "SELECT tool_name FROM registered_tools WHERE id = ?", (tool_id,)
    ).fetchone()
    if not row:
        raise HTTPException(404, "Tool not found")
    tool_name = row["tool_name"]
    conn.execute("DELETE FROM registered_tools WHERE id = ?", (tool_id,))
    conn.commit()
    conn.close()
    AI_TOOLS_REGISTRY.pop(tool_name, None)
    log_tool_change(tool_name, "deleted")
    return {"status": "ok"}


@router.get("/api/tool-change-log")
async def get_tool_change_log(page: int = 1, per_page: int = 30, tool_name: str = ""):
    conn = get_db()
    filters = [
        "action != 'file_deleted'",
        "NOT (action = 'updated' AND (changed_fields IS NULL OR changed_fields = '{}' OR changed_fields = ''))",
    ]
    args = []
    if tool_name.strip():
        filters.append("tool_name LIKE ?")
        args.append(f"%{tool_name.strip()}%")
    where = "WHERE " + " AND ".join(filters)
    base = f"FROM tool_change_log {where}"
    total = conn.execute(f"SELECT COUNT(*) as c {base}", args).fetchone()["c"]
    offset = (page - 1) * per_page
    rows = conn.execute(
        f"SELECT * {base} ORDER BY created_at DESC OFFSET ? ROWS FETCH NEXT ? ROWS ONLY",
        args + [offset, per_page],
    ).fetchall()
    conn.close()
    return {
        "total":    total,
        "page":     page,
        "per_page": per_page,
        "logs":     [dict(r) for r in rows],
    }


@router.delete("/api/tool-docs/file")
async def delete_tool_doc_file(filename: str, tool_name: str = ""):
    try:
        results = tool_knowledge_collection.get(
            where={"source_file": filename}, include=[]
        )
        ids = results.get("ids", [])
        if ids:
            tool_knowledge_collection.delete(ids=ids)
        log_tool_change(
            tool_name=tool_name or filename,
            action="updated",
            changed_fields={"file_removed": filename, "chunks_removed": len(ids)},
        )
        return {"status": "ok", "deleted_chunks": len(ids), "filename": filename}
    except Exception as e:
        raise HTTPException(500, str(e))


@router.post("/api/tool-docs/upload/{tool_name}")
async def upload_tool_docs_for_tool(tool_name: str, files: List[UploadFile] = File(...)):
    processed_files = []
    total_chunks    = 0

    for f in files:
        content  = await f.read()
        filename = f.filename or "unknown"
        try:
            summary = ingest_tool_document_direct(content, filename, tool_name)
            count   = summary.get(tool_name, 0)
            total_chunks += count
            processed_files.append({"filename": filename, "chunks": count})
        except Exception as exc:
            processed_files.append({"filename": filename, "error": str(exc)})

    filenames = [f["filename"] for f in processed_files if "error" not in f]
    if filenames:
        log_tool_change(
            tool_name=tool_name,
            action="updated",
            changed_fields={
                "files_uploaded": ", ".join(filenames),
                "chunks_added":   total_chunks,
            },
        )

    return {
        "status":          "ok",
        "tool_name":       tool_name,
        "files_processed": len(processed_files),
        "files":           processed_files,
        "total_chunks":    total_chunks,
    }


@router.post("/api/tool-docs/upload")
async def upload_tool_docs(files: List[UploadFile] = File(...)):
    tool_names = list(AI_TOOLS_REGISTRY.keys())
    if not tool_names:
        raise HTTPException(400, "No tools in registry. Upload your Excel first so the classifier knows what tools exist.")

    overall_summary = {}
    processed_files = []

    for f in files:
        content  = await f.read()
        filename = f.filename or "unknown"
        try:
            file_summary = ingest_tool_document(content, filename, tool_names)
            for tool, count in file_summary.items():
                overall_summary[tool] = overall_summary.get(tool, 0) + count
            processed_files.append({
                "filename":    filename,
                "chunks":      sum(file_summary.values()),
                "classified":  {k: v for k, v in file_summary.items() if k != "unclassified"},
                "unclassified": file_summary.get("unclassified", 0),
            })
        except Exception as exc:
            processed_files.append({"filename": filename, "error": str(exc)})

    return {
        "status":          "ok",
        "files_processed": len(processed_files),
        "files":           processed_files,
        "tool_summary":    {k: v for k, v in overall_summary.items() if k != "unclassified"},
        "unclassified_chunks": overall_summary.get("unclassified", 0),
    }


@router.get("/api/tool-docs/status")
async def tool_docs_status():
    return {"status": get_tool_knowledge_status()}


@router.delete("/api/tool-docs/clear")
async def clear_tool_docs():
    try:
        all_ids = tool_knowledge_collection.get()["ids"]
        if all_ids:
            tool_knowledge_collection.delete(ids=all_ids)
        return {"status": "ok", "deleted": len(all_ids)}
    except Exception as e:
        raise HTTPException(500, str(e))


@router.delete("/api/tool-docs/clear/{tool_name}")
async def clear_tool_docs_for_tool(tool_name: str):
    try:
        results = tool_knowledge_collection.get(
            where={"tool_name": tool_name}, include=[]
        )
        ids = results.get("ids", [])
        if ids:
            tool_knowledge_collection.delete(ids=ids)
        return {"status": "ok", "deleted": len(ids), "tool_name": tool_name}
    except Exception as e:
        raise HTTPException(500, str(e))


@router.post("/api/upload-tools-registry")
async def upload_tools_registry(file: UploadFile = File(...)):
    filename = file.filename or ""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext not in ("xlsx", "xlsm", "xls"):
        raise HTTPException(400, "Only Excel files (.xlsx, .xlsm, .xls) are supported")

    content = await file.read()

    try:
        reload_tools_registry(excel_bytes=content)
    except ValueError as e:
        raise HTTPException(400, str(e))
    except Exception as e:
        raise HTTPException(400, f"Could not read Excel file: {str(e)}")

    if not AI_TOOLS_REGISTRY:
        raise HTTPException(400, "File was read but no tools were found. Ensure your sheet has a tool name column and at least one data row.")

    return {"status": "ok", "tools_loaded": len(AI_TOOLS_REGISTRY)}
