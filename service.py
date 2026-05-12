import os
import json
import uuid
import sqlite3
import threading
from datetime import datetime
from typing import List
import pandas as pd

import chromadb
from chromadb.utils import embedding_functions
from langgraph.graph import StateGraph, END

from schemas import OrchestratorState

# ══════════════════════════════════════════════════════════════════════════════
# AZURE OPENAI CLIENT
# ══════════════════════════════════════════════════════════════════════════════
try:
    from openai import AzureOpenAI

    _azure_client = AzureOpenAI(
        api_key=os.getenv("AZURE_OPENAI_API_KEY", ""),
        azure_endpoint=os.getenv("AZURE_OPENAI_BASE_URL", ""),
        api_version=os.getenv("AZURE_OPENAI_API_VERSION", "2024-02-01"),
    )
    _AZURE_DEPLOYMENT = os.getenv("AZURE_OPENAI_DEPLOYMENT", "")
    HAS_AZURE = bool(
        os.getenv("AZURE_OPENAI_API_KEY") and
        os.getenv("AZURE_OPENAI_BASE_URL") and
        _AZURE_DEPLOYMENT
    )
except Exception:
    _azure_client = None
    _AZURE_DEPLOYMENT = ""
    HAS_AZURE = False

SYSTEM_VERSION = "2.0"

# ══════════════════════════════════════════════════════════════════════════════
# TAVILY SEARCH CLIENT  (optional — enriches tool profiles with real-world data)
# Set TAVILY_API_KEY in .env to enable. Works without it — just skips enrichment.
# ══════════════════════════════════════════════════════════════════════════════
try:
    from tavily import TavilyClient as _TavilyClient
    _tavily    = _TavilyClient(api_key=os.getenv("TAVILY_API_KEY", ""))
    HAS_TAVILY = bool(os.getenv("TAVILY_API_KEY"))
except Exception:
    _tavily    = None
    HAS_TAVILY = False

# ── In-memory enrichment cache ────────────────────────────────────────────────
# Key: tool_name → {"summary": str, "fetched_at": str}
# Cleared automatically when the Excel registry is reloaded (new upload).
# This means we only call Tavily once per tool per session, not on every request.
_TOOL_ENRICHMENT_CACHE: dict = {}

# ══════════════════════════════════════════════════════════════════════════════
# AI TOOLS CAPABILITY REGISTRY (Loaded from Excel)
# ══════════════════════════════════════════════════════════════════════════════

def _split_list(val) -> list:
    """Split a comma- or semicolon-separated Excel cell value into a clean list."""
    if val is None:
        return []
    s = str(val).strip()
    if not s or s.lower() in ("nan", "none", ""):
        return []
    sep = ";" if ";" in s else ","
    return [x.strip() for x in s.split(sep) if x.strip()]


# ── Live registry — mutated in-place when user uploads a new Excel ────────────
AI_TOOLS_REGISTRY: dict = {}


# ══════════════════════════════════════════════════════════════════════════════
# EXCEL LOADER  — accepts any filename, any sheet name, any column names.
#
# Strategy:
#   1. Auto-detect the correct sheet (name matching + column sniffing).
#   2. Find the one column that identifies the tool name (required).
#   3. Capture EVERY other column as-is into raw_data — zero data loss.
#   4. Derive the 7 fields that downstream code reads (.get("description") etc.)
#      via a priority-ordered alias list.  If a column maps, great; if not,
#      the field stays "" / [] — callers already guard with .get("x", default).
#
# Result per registry entry:
#   {
#     # ── derived fields (consumed by routing & prompt nodes) ──
#     "description":  str,
#     "category":     str,
#     "url":          str,
#     "icon":         str,
#     "search_query": str,
#     "is_internal":  bool,
#     "best_for":     [str, ...],
#     "not_for":      [str, ...],
#     "roles":        [str, ...],
#     # ── full raw capture (every column, every value) ──
#     "raw_data":     { "col_name": value, ... }
#   }
# ══════════════════════════════════════════════════════════════════════════════

# Alias lists for the 9 derived fields — ordered by preference.
# Any column whose normalised name appears in the list maps to that field.
_FIELD_ALIASES: dict = {
    "tool_name":            ["tool_name", "name", "tool", "ai_tool", "toolname", "tool_names"],
    "description":          ["description", "desc", "primary_job", "summary", "overview", "about", "what_it_does"],
    "category":             ["category", "cat", "tool_category", "tool_family", "family", "type", "group"],
    "url":                  ["url", "url_or_note", "link", "tool_url", "access_url", "url_link", "endpoint"],
    "best_for":             ["best_for", "use_when", "use_case", "good_for", "when_to_use", "ideal_for"],
    "strong_signals":       ["strong_signals", "signals", "keywords", "match_keywords", "trigger_keywords"],
    "not_for":              ["not_for", "avoid_when", "avoid", "do_not_use", "dont_use", "not_suitable"],
    "weak_signals":         ["weak_signals", "weak", "secondary_signals"],
    "roles":                ["roles", "role", "target_roles", "audience", "user_roles", "who", "for_roles"],
    "icon":                 ["icon", "emoji", "symbol"],
    "search_query":         ["search_query", "search", "query"],
    "is_internal":          ["is_internal", "internal", "internal_only", "allows_client_data"],
    "output_type":          ["output_type", "output", "produces", "delivers", "tool_output"],
    "output_type_keywords": ["output_type_keywords", "output_keywords", "output_signals"],
}


def _norm(s: str) -> str:
    """Normalise a column name for alias matching."""
    return s.lower().strip().replace(" ", "_").replace("-", "_")


def _safe_val(val) -> str:
    """Return a clean string from any cell value, '' on NaN/None."""
    if val is None:
        return ""
    if isinstance(val, float) and pd.isna(val):
        return ""
    return str(val).strip()


def _find_tool_sheet(xl: pd.ExcelFile) -> str:
    """
    Auto-detect the tools sheet.
    Pass 1 — name match against known canonical names.
    Pass 2 — scan each sheet for a tool_name-like column.
    """
    canonical = {
        "ai_tools_registry", "registry_all_tools", "tools_registry",
        "tool_registry", "ai_tools", "registry", "tools", "tool_list",
    }
    for sheet in xl.sheet_names:
        if _norm(sheet) in canonical:
            return sheet

    for sheet in xl.sheet_names:
        try:
            df_peek = xl.parse(sheet, nrows=1)
            normed  = [_norm(c) for c in df_peek.columns]
            if any(
                alias in normed
                for alias in _FIELD_ALIASES["tool_name"]
            ):
                return sheet
        except Exception:
            continue

    raise ValueError(
        f"No tools registry sheet found. Sheets in this file: {xl.sheet_names}. "
        "Ensure one sheet has a column named 'tool_name' (or similar)."
    )


def _resolve_aliases(col_lookup: dict) -> dict:
    """
    col_lookup: { normalised_col_name → original_col_name }
    Returns: { internal_field → original_col_name | None }
    """
    resolved = {}
    for field, aliases in _FIELD_ALIASES.items():
        resolved[field] = next(
            (col_lookup[a] for a in aliases if a in col_lookup),
            None
        )
    return resolved


def _load_from_bytes(excel_bytes: bytes) -> dict:
    """
    Core loader. Accepts raw Excel bytes, returns the populated registry dict.
    Works with any filename, any sheet name, any column naming convention.
    """
    import io

    try:
        xl = pd.ExcelFile(io.BytesIO(excel_bytes), engine="openpyxl")
    except Exception as e:
        raise ValueError(f"Could not open Excel file: {e}")

    sheet = _find_tool_sheet(xl)

    try:
        df = xl.parse(sheet)
    except Exception as e:
        raise ValueError(f"Could not parse sheet '{sheet}': {e}")

    # Drop rows and columns that are entirely empty
    df = df.dropna(how="all").reset_index(drop=True)
    df.columns = [str(c) for c in df.columns]   # ensure all column names are strings

    # Build lookup: normalised name → original name
    col_lookup = {_norm(c): c for c in df.columns}

    # Resolve aliases → actual column names (or None if not present)
    alias_map = _resolve_aliases(col_lookup)

    if not alias_map["tool_name"]:
        raise ValueError(
            f"Sheet '{sheet}' has no recognisable tool-name column. "
            f"Columns found: {list(df.columns)}. "
            f"Rename the tool name column to 'tool_name'."
        )

    registry = {}

    for _, row in df.iterrows():

        # ── 1. Identify the tool name (the registry key) ──────────────────────
        tool_name = _safe_val(row[alias_map["tool_name"]])
        if not tool_name or tool_name.lower() in ("nan", "none", "tool_name", "tool name"):
            continue

        # ── 2. Capture ALL columns into raw_data ──────────────────────────────
        raw_data = {}
        for orig_col in df.columns:
            v = row[orig_col]
            # Store as native Python type (str, int, float, bool) — JSON-safe
            if v is None or (isinstance(v, float) and pd.isna(v)):
                raw_data[orig_col] = None
            elif isinstance(v, (int, float, bool)):
                raw_data[orig_col] = v
            else:
                raw_data[orig_col] = str(v).strip()

        # ── 3. Derive the 9 standard fields via alias resolution ───────────────
        def _get(field: str) -> str:
            col = alias_map.get(field)
            return _safe_val(row[col]) if col else ""

        # is_internal needs special handling:
        # - if mapped to "allows_client_data" → invert ("Yes" = client data OK = NOT internal)
        # - otherwise → truthy string check
        is_internal = False
        is_internal_col = alias_map.get("is_internal")
        if is_internal_col:
            raw_flag = _safe_val(row[is_internal_col]).lower()
            if _norm(is_internal_col) == "allows_client_data":
                is_internal = raw_flag not in ("yes", "true", "1")
            else:
                is_internal = raw_flag in ("true", "yes", "1", "internal")

        registry[tool_name] = {
            # Derived fields — consumed by routing nodes & prompt builders
            "description":          _get("description"),
            "category":             _get("category"),
            "url":                  _get("url"),
            "icon":                 _get("icon"),
            "search_query":         _get("search_query"),
            "is_internal":          is_internal,
            "best_for":             _split_list(_get("best_for")),
            "strong_signals":       _split_list(_get("strong_signals")),
            "not_for":              _split_list(_get("not_for")),
            "weak_signals":         _split_list(_get("weak_signals")),
            "roles":                _split_list(_get("roles")),
            "output_type":          _get("output_type"),
            "output_type_keywords": _split_list(_get("output_type_keywords")),
            # Full raw capture — every column, every value, nothing dropped
            "raw_data":             raw_data,
        }

    if not registry:
        raise ValueError(
            f"Sheet '{sheet}' was parsed but contains no valid tool rows. "
            "Ensure the sheet has at least one data row with a tool name."
        )

    return registry


def load_tools_registry_from_excel(
    excel_path: str = "AI_TOOLS.xlsx",
    sheet_name: str = None,          # ignored — kept for backward-compat call signature
) -> dict:
    """Disk-path wrapper around _load_from_bytes. Used on startup."""
    with open(excel_path, "rb") as fh:
        return _load_from_bytes(fh.read())


def reload_tools_registry(excel_bytes: bytes = None,
                           excel_path:  str   = "AI_TOOLS.xlsx"):
    """
    (Re)load the registry from either uploaded bytes or a disk path.
    Accepts any Excel filename, any sheet name, any column naming convention.
    Mutates AI_TOOLS_REGISTRY in-place so every node sees the update instantly.
    DB-registered tools are merged in and take precedence over Excel entries
    with the same name.
    """
    global AI_TOOLS_REGISTRY

    if excel_bytes:
        new = _load_from_bytes(excel_bytes)
    else:
        new = load_tools_registry_from_excel(excel_path=excel_path)

    if not new:
        raise ValueError(
            "Excel file was parsed but no tools were found. "
            "Ensure your sheet has a 'tool_name' column and at least one data row."
        )

    AI_TOOLS_REGISTRY.clear()
    AI_TOOLS_REGISTRY.update(new)
    _merge_db_tools_into_registry()
    _TOOL_ENRICHMENT_CACHE.clear()


def _merge_db_tools_into_registry():
    """Load all tools from registered_tools table and merge into AI_TOOLS_REGISTRY."""
    try:
        conn = get_db()
        rows = conn.execute("SELECT * FROM registered_tools").fetchall()
        conn.close()
        for row in rows:
            r = dict(row)
            tool_name = r["tool_name"]
            AI_TOOLS_REGISTRY[tool_name] = {
                "description":         r.get("description", ""),
                "category":            r.get("category", ""),
                "url":                 r.get("url", ""),
                "icon":                r.get("icon", "🤖"),
                "best_for":            json.loads(r.get("best_for", "[]") or "[]"),
                "strong_signals":      json.loads(r.get("strong_signals", "[]") or "[]"),
                "weak_signals":        json.loads(r.get("weak_signals", "[]") or "[]"),
                "not_for":             json.loads(r.get("not_for", "[]") or "[]"),
                "roles":               json.loads(r.get("roles", "[]") or "[]"),
                "output_type":         r.get("output_type", ""),
                "output_type_keywords": [],
                "is_internal":         bool(r.get("is_internal", 0)),
                "search_query":        "",
                "raw_data":            json.loads(r.get("raw_data", "{}") or "{}"),
                "_source":             "db",
            }
    except Exception:
        pass


# Load from disk on startup
try:
    AI_TOOLS_REGISTRY.update(load_tools_registry_from_excel())
except Exception as e:
    import warnings
    warnings.warn(
        f"[AI_TOOLS_REGISTRY] Failed to load Excel on startup: {e}. "
        "Upload a registry via the UI header dropdown before using the tool recommender.",
        RuntimeWarning,
        stacklevel=1,
    )

# Always merge DB-registered tools on startup (even if Excel failed)
try:
    _merge_db_tools_into_registry()
except Exception:
    pass




# ══════════════════════════════════════════════════════════════════════════════
# SCENARIO LIBRARY  — loaded from AI_Navigator_Scenario_Library_Refined.xlsx
# Structure per row: Mega-Group, Category, Scenario Title, Persona / Role, Scenarios
# Stored as:
#   SCENARIO_LIBRARY = [
#     { "mega_group": str, "category": str, "title": str, "persona": str, "scenario": str },
#     ...
#   ]
# ══════════════════════════════════════════════════════════════════════════════

SCENARIO_LIBRARY: list = []


def _load_scenario_library_from_bytes(excel_bytes: bytes) -> list:
    """Parse the scenario library Excel and return a list of scenario dicts."""
    import io as _io
    try:
        xl = pd.ExcelFile(_io.BytesIO(excel_bytes), engine="openpyxl")
    except Exception as e:
        raise ValueError(f"Could not open Excel file: {e}")

    # Find the correct sheet — prefer 'Scenario Library', otherwise first sheet
    target_sheet = xl.sheet_names[0]
    for sh in xl.sheet_names:
        if _norm(sh) in ("scenario_library", "scenarios", "scenario", "library"):
            target_sheet = sh
            break

    try:
        df = xl.parse(target_sheet)
    except Exception as e:
        raise ValueError(f"Could not parse sheet '{target_sheet}': {e}")

    df = df.dropna(how="all").reset_index(drop=True)
    df.columns = [str(c) for c in df.columns]

    # Normalised column lookup
    col_lookup = {_norm(c): c for c in df.columns}

    # Map expected fields using flexible aliases
    _scenario_aliases = {
        "mega_group": ["mega_group", "mega-group", "mega group", "group", "mega"],
        "category":   ["category", "cat", "sub_group", "sub-group", "subgroup"],
        "phase":      ["activate_phase", "activate phase", "phase", "sap_phase", "sap phase"],
        "title":      ["scenario_title", "title", "scenario title", "name"],
        "persona":    ["persona_/_role", "persona", "role", "persona_role", "persona / role"],
        "scenario":   ["scenarios", "scenario", "description", "prompt", "task", "body"],
    }

    def _find_col(key):
        for alias in _scenario_aliases[key]:
            n = _norm(alias)
            if n in col_lookup:
                return col_lookup[n]
        return None

    col_mega     = _find_col("mega_group")
    col_category = _find_col("category")
    col_phase    = _find_col("phase")
    col_title    = _find_col("title")
    col_persona  = _find_col("persona")
    col_scenario = _find_col("scenario")

    if not col_title and not col_scenario:
        raise ValueError(
            f"Sheet '{target_sheet}' has no recognisable title/scenario column. "
            f"Columns found: {list(df.columns)}"
        )

    scenarios = []
    last_mega = ""
    last_cat  = ""

    for _, row in df.iterrows():
        mega     = _safe_val(row[col_mega])     if col_mega     else ""
        category = _safe_val(row[col_category]) if col_category else ""
        phase    = _safe_val(row[col_phase])    if col_phase    else ""
        title    = _safe_val(row[col_title])    if col_title    else ""
        persona  = _safe_val(row[col_persona])  if col_persona  else ""
        scenario = _safe_val(row[col_scenario]) if col_scenario else ""

        # Carry forward mega-group / category for merged cells
        if mega:
            last_mega = mega
        else:
            mega = last_mega

        if category:
            last_cat = category
        else:
            category = last_cat

        # Skip header-like rows or empty rows
        if not title and not scenario:
            continue
        if title.lower() in ("scenario title", "title", "name"):
            continue

        phase_clean = phase if phase and phase not in ("-", "—", "–") else ""
        scenarios.append({
            "mega_group": mega,
            "category":   category,
            "phase":      phase_clean,
            "title":      title,
            "persona":    persona,
            "scenario":   scenario,
        })

    return scenarios


def reload_scenario_library(excel_bytes: bytes = None,
                             excel_path: str = "AI_Navigator_Scenario_Library_Refined.xlsx"):
    """
    (Re)load the scenario library from uploaded bytes or a disk path.
    Mutates SCENARIO_LIBRARY in-place.
    """
    global SCENARIO_LIBRARY

    if excel_bytes:
        new = _load_scenario_library_from_bytes(excel_bytes)
    else:
        with open(excel_path, "rb") as fh:
            new = _load_scenario_library_from_bytes(fh.read())

    SCENARIO_LIBRARY.clear()
    SCENARIO_LIBRARY.extend(new)


# Load from disk on startup
try:
    reload_scenario_library()
except Exception as _e:
    import warnings
    warnings.warn(
        f"[SCENARIO_LIBRARY] Failed to load Excel on startup: {_e}. "
        "Upload via the Scenario Library upload button in the UI.",
        RuntimeWarning,
        stacklevel=1,
    )


# ══════════════════════════════════════════════════════════════════════════════
# TAVILY TOOL ENRICHMENT
# Fetches real-world capability summaries for each tool in the registry.
# Results are cached in _TOOL_ENRICHMENT_CACHE for the session lifetime.
# Called once per tool, not on every request.
# ══════════════════════════════════════════════════════════════════════════════

def _build_search_query(tool_name: str, info: dict) -> str:
    """
    Build the best possible Tavily search query for a tool.
    Priority:
      1. Explicit search_query from Excel (admin-curated — most precise)
      2. Auto-generated from tool_name + category
    """
    explicit = info.get("search_query", "").strip()
    if explicit:
        return explicit
    category = info.get("category", "").strip()
    if category:
        return f"{tool_name} {category} features capabilities use cases enterprise"
    return f"{tool_name} AI tool features capabilities what it does"


def _enrich_single_tool(tool_name: str, info: dict) -> str:
    """
    Fetch a capability summary for one tool via Tavily.
    Hard timeout of 8 seconds so it never blocks a request indefinitely.
    Falls back gracefully to the Excel description if Tavily fails or times out.
    """
    if not HAS_TAVILY or _tavily is None:
        return info.get("description", "")

    query = _build_search_query(tool_name, info)
    result_holder = [None]

    def _do_search():
        try:
            result_holder[0] = _tavily.search(
                query=query,
                search_depth="basic",
                max_results=3,
                include_answer=True,
            )
        except Exception:
            result_holder[0] = {}

    t = threading.Thread(target=_do_search, daemon=True)
    t.start()
    t.join(timeout=8)

    result = result_holder[0] or {}
    answer = (result.get("answer") or "").strip()
    if answer:
        return answer[:600]

    snippets = [
        r.get("content", "")[:150]
        for r in (result.get("results") or [])[:3]
        if r.get("content")
    ]
    combined = " ".join(snippets).strip()
    return combined[:600] if combined else info.get("description", "")


def enrich_tools_registry() -> None:
    """
    Enrich all tools in AI_TOOLS_REGISTRY that are not yet in the cache.
    Runs in a background daemon thread so it never blocks any request.
    Safe to call multiple times — skips already-cached tools.
    """
    if not AI_TOOLS_REGISTRY:
        return

    new_tools = [
        name for name in AI_TOOLS_REGISTRY
        if name not in _TOOL_ENRICHMENT_CACHE
    ]

    if not new_tools:
        return

    def _run():
        for name in new_tools:
            if name in _TOOL_ENRICHMENT_CACHE:
                continue
            info    = AI_TOOLS_REGISTRY.get(name, {})
            summary = _enrich_single_tool(name, info)
            _TOOL_ENRICHMENT_CACHE[name] = {
                "summary":    summary,
                "fetched_at": datetime.utcnow().isoformat(),
            }

    threading.Thread(target=_run, daemon=True).start()


# Kick off background enrichment after function is defined — never blocks requests
threading.Thread(target=enrich_tools_registry, daemon=True).start()


def _get_enriched_summary(tool_name: str) -> str:
    """
    Returns the enriched summary for a tool.
    Falls back to the Excel description if enrichment is not available.
    """
    cached = _TOOL_ENRICHMENT_CACHE.get(tool_name)
    if cached and cached.get("summary"):
        return cached["summary"]
    info = AI_TOOLS_REGISTRY.get(tool_name, {})
    return info.get("description", "")


# ══════════════════════════════════════════════════════════════════════════════
# CHROMADB SETUP
# ══════════════════════════════════════════════════════════════════════════════
chroma_client = chromadb.PersistentClient(path="./chroma_db")
ef = embedding_functions.DefaultEmbeddingFunction()
policy_collection = chroma_client.get_or_create_collection(
    name="company_policies", embedding_function=ef
)
tool_knowledge_collection = chroma_client.get_or_create_collection(
    name="tool_knowledge", embedding_function=ef
)


# ══════════════════════════════════════════════════════════════════════════════
# TOOL KNOWLEDGE BASE — chunk, classify and store tool documents in ChromaDB
# ══════════════════════════════════════════════════════════════════════════════

def _chunk_text(text: str, chunk_size: int = 400, overlap: int = 60) -> list:
    """Split text into overlapping word-level chunks."""
    words  = text.split()
    chunks = []
    start  = 0
    while start < len(words):
        end = start + chunk_size
        chunks.append(" ".join(words[start:end]))
        start += chunk_size - overlap
    return [c.strip() for c in chunks if len(c.strip()) > 60]


def _classify_chunk_to_tool(chunk: str, tool_names: list) -> str:
    """
    Ask the LLM which tool from the registry this text chunk describes.
    Returns the matching tool name or 'unclassified'.
    """
    if not HAS_AZURE or not _azure_client or not tool_names:
        return "unclassified"

    names_list = "\n".join(f"- {n}" for n in tool_names)
    system_msg = (
        "You are a document classifier. Given a text chunk from a tool document, "
        "identify which AI tool it describes from the provided list. "
        "Reply with ONLY the exact tool name from the list, or 'unclassified' if unsure."
    )
    user_msg = (
        f"Tool names to choose from:\n{names_list}\n\n"
        f"Text chunk:\n{chunk[:800]}\n\n"
        "Which tool does this text describe? Reply with only the exact tool name or 'unclassified'."
    )
    try:
        result = call_llm(system_msg, user_msg, max_tokens=30, temperature=0.0)
        result = result.strip().strip('"').strip("'")
        if result in tool_names:
            return result
        # fuzzy match — case-insensitive
        result_lower = result.lower()
        for name in tool_names:
            if name.lower() == result_lower or name.lower() in result_lower:
                return name
        return "unclassified"
    except Exception:
        return "unclassified"


def ingest_tool_document_direct(file_bytes: bytes, filename: str, tool_name: str) -> dict:
    """
    Extract text from a document, chunk it, and store every chunk directly
    under the given tool_name — no LLM classifier involved.

    Returns a summary dict: { tool_name: chunk_count }
    """
    from routes import _extract_pdf_text, _extract_docx_text

    fname = filename.lower()
    text  = ""
    try:
        if fname.endswith(".pdf"):
            text = _extract_pdf_text(file_bytes)
        elif fname.endswith(".docx"):
            text = _extract_docx_text(file_bytes)
        elif fname.endswith(".pptx"):
            try:
                from pptx import Presentation
                from io import BytesIO
                prs = Presentation(BytesIO(file_bytes))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text += shape.text.strip() + "\n"
            except Exception:
                text = file_bytes.decode("utf-8", errors="ignore")
        else:
            text = file_bytes.decode("utf-8", errors="ignore")
    except Exception:
        text = ""

    if not text.strip():
        return {}

    chunks = _chunk_text(text)
    count  = 0

    for idx, chunk in enumerate(chunks):
        chunk_id = f"{tool_name}__{filename}__chunk_{idx}"
        try:
            tool_knowledge_collection.upsert(
                ids       = [chunk_id],
                documents = [chunk],
                metadatas = [{
                    "tool_name":   tool_name,
                    "source_file": filename,
                    "chunk_index": idx,
                }],
            )
            count += 1
        except Exception:
            pass

    return {tool_name: count}


def ingest_tool_document(file_bytes: bytes, filename: str, tool_names: list) -> dict:
    """
    Extract text from a document, chunk it, classify each chunk to a tool,
    and store everything in the tool_knowledge ChromaDB collection.

    Returns a summary dict: { tool_name: chunk_count, ... }
    """
    from routes import _extract_pdf_text, _extract_docx_text

    fname = filename.lower()
    text  = ""
    try:
        if fname.endswith(".pdf"):
            text = _extract_pdf_text(file_bytes)
        elif fname.endswith(".docx"):
            text = _extract_docx_text(file_bytes)
        elif fname.endswith(".pptx"):
            try:
                from pptx import Presentation
                from io import BytesIO
                prs = Presentation(BytesIO(file_bytes))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text += shape.text.strip() + "\n"
            except Exception:
                text = file_bytes.decode("utf-8", errors="ignore")
        else:
            text = file_bytes.decode("utf-8", errors="ignore")
    except Exception:
        text = ""

    if not text.strip():
        return {}

    chunks  = _chunk_text(text)
    summary = {}

    for idx, chunk in enumerate(chunks):
        tool_name = _classify_chunk_to_tool(chunk, tool_names)
        chunk_id  = f"{filename}__chunk_{idx}"

        try:
            tool_knowledge_collection.upsert(
                ids        = [chunk_id],
                documents  = [chunk],
                metadatas  = [{
                    "tool_name":   tool_name,
                    "source_file": filename,
                    "chunk_index": idx,
                }],
            )
        except Exception:
            pass

        summary[tool_name] = summary.get(tool_name, 0) + 1

    return summary


def query_tool_knowledge(user_input: str, tool_names: list, n_results: int = 5) -> dict:
    """
    Semantic search in tool_knowledge for chunks relevant to the user's task.
    Returns { tool_name: [chunk_text, ...] } — only for tools in tool_names.
    """
    if not tool_names:
        return {}
    try:
        results = tool_knowledge_collection.query(
            query_texts = [user_input],
            n_results   = min(n_results, 10),
        )
    except Exception:
        return {}

    docs      = results.get("documents",  [[]])[0]
    metas     = results.get("metadatas",  [[]])[0]
    knowledge = {}

    for doc, meta in zip(docs, metas):
        t = meta.get("tool_name", "unclassified")
        if t in tool_names:
            knowledge.setdefault(t, []).append(doc)

    return knowledge


def get_tool_knowledge_status() -> list:
    """
    Return a summary of what's stored in tool_knowledge:
    [{ tool_name, chunk_count, source_files }, ...]
    """
    try:
        all_items = tool_knowledge_collection.get(include=["metadatas"])
        metas     = all_items.get("metadatas", [])
    except Exception:
        return []

    status = {}
    for m in metas:
        t     = m.get("tool_name", "unclassified")
        fname = m.get("source_file", "")
        if t not in status:
            status[t] = {"tool_name": t, "chunk_count": 0, "source_files": set()}
        status[t]["chunk_count"]  += 1
        status[t]["source_files"].add(fname)

    return [
        {**v, "source_files": sorted(v["source_files"])}
        for v in sorted(status.values(), key=lambda x: x["chunk_count"], reverse=True)
    ]


# ══════════════════════════════════════════════════════════════════════════════
# SQLITE DATABASE
# ══════════════════════════════════════════════════════════════════════════════
DB_PATH = "./orchestrator.db"


def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def log_tool_change(tool_name: str, action: str, changed_fields: dict = None, note: str = ""):
    """Write an entry to the tool_change_log table."""
    try:
        conn = get_db()
        conn.execute(
            "INSERT INTO tool_change_log (id, tool_name, action, changed_fields, changed_by, note, created_at) "
            "VALUES (?,?,?,?,?,?,?)",
            (
                str(uuid.uuid4()),
                tool_name,
                action,
                json.dumps(changed_fields or {}),
                "admin",
                note,
                datetime.utcnow().isoformat(),
            ),
        )
        conn.commit()
        conn.close()
    except Exception:
        pass


def init_db():
    conn = get_db()
    conn.executescript("""
        CREATE TABLE IF NOT EXISTS audit_log (
            id TEXT PRIMARY KEY,
            created_at TEXT,
            raw_input TEXT,
            intent TEXT,
            industry TEXT,
            recommended_tool TEXT,
            tool_reason TEXT,
            tool_confidence TEXT,
            policy_flags TEXT,
            retrieved_policies TEXT,
            final_prompt TEXT,
            prompt_version TEXT,
            model_used TEXT,
            output TEXT,
            token_estimate INTEGER,
            system_version TEXT,
            policy_blocked INTEGER DEFAULT 0,
            policy_summary TEXT DEFAULT ''
        );
        -- Add columns to existing DB if upgrading (safe: ignored if already present)
        -- SQLite ignores "duplicate column" errors via the try/except in init_db

        CREATE TABLE IF NOT EXISTS feedback (
            id TEXT PRIMARY KEY,
            audit_id TEXT,
            email TEXT DEFAULT '',
            rating INTEGER,
            comment TEXT,
            issue_type TEXT,
            created_at TEXT,
            source TEXT DEFAULT 'form'
        );
        CREATE TABLE IF NOT EXISTS prompt_versions (
            id TEXT PRIMARY KEY,
            version TEXT NOT NULL,
            intent TEXT,
            industry TEXT,
            template TEXT NOT NULL,
            change_note TEXT,
            created_at TEXT,
            created_by TEXT DEFAULT 'system'
        );
        CREATE TABLE IF NOT EXISTS registered_tools (
            id TEXT PRIMARY KEY,
            tool_name TEXT NOT NULL UNIQUE,
            description TEXT DEFAULT '',
            category TEXT DEFAULT '',
            url TEXT DEFAULT '',
            icon TEXT DEFAULT '🤖',
            best_for TEXT DEFAULT '[]',
            strong_signals TEXT DEFAULT '[]',
            weak_signals TEXT DEFAULT '[]',
            not_for TEXT DEFAULT '[]',
            roles TEXT DEFAULT '[]',
            output_type TEXT DEFAULT '',
            is_internal INTEGER DEFAULT 0,
            raw_data TEXT DEFAULT '{}',
            created_at TEXT,
            updated_at TEXT
        );
        CREATE TABLE IF NOT EXISTS tool_change_log (
            id TEXT PRIMARY KEY,
            tool_name TEXT NOT NULL,
            action TEXT NOT NULL,
            changed_fields TEXT DEFAULT '{}',
            changed_by TEXT DEFAULT 'admin',
            note TEXT DEFAULT '',
            created_at TEXT
        );
        CREATE TABLE IF NOT EXISTS scenario_suggestions (
            id TEXT PRIMARY KEY,
            title TEXT NOT NULL,
            mega_group TEXT NOT NULL,
            category TEXT DEFAULT '',
            persona TEXT DEFAULT '',
            activate_phase TEXT DEFAULT '',
            scenario TEXT NOT NULL,
            submitted_by TEXT DEFAULT '',
            submitted_at TEXT NOT NULL,
            status TEXT DEFAULT 'pending',
            admin_note TEXT DEFAULT '',
            reviewed_at TEXT DEFAULT ''
        );
    """)
    # Migrate existing DB — add new columns if they don't exist yet
    for col, definition in [("policy_blocked", "INTEGER DEFAULT 0"),
                             ("policy_summary", "TEXT DEFAULT ''"),
                             ("role",           "TEXT DEFAULT 'general'"),
                             ("user_email",     "TEXT DEFAULT ''")]:
        try:
            conn.execute(f"ALTER TABLE audit_log ADD COLUMN {col} {definition}")
            conn.commit()
        except Exception:
            pass  # Column already exists
    for col, definition in [("email", "TEXT DEFAULT ''"), ("source", "TEXT DEFAULT 'form'")]:
        try:
            conn.execute(f"ALTER TABLE feedback ADD COLUMN {col} {definition}")
            conn.commit()
        except Exception:
            pass  # Column already exists

    count = conn.execute("SELECT COUNT(*) as c FROM prompt_versions").fetchone()["c"]
    if count == 0:
        conn.execute(
            "INSERT INTO prompt_versions VALUES (?,?,?,?,?,?,?,?)",
            (
                str(uuid.uuid4()), "1.0", "general", "general",
                "## ROLE\nYou are an expert {industry} professional specializing in {intent} tasks.\n\n"
                "## CONTEXT\nUser Request: {user_input}\nIndustry: {industry} | Task Type: {intent}\nTarget Tool: {tool}\n\n"
                "## OBJECTIVE\nProduce a high-quality, professional {intent} that directly addresses the user's need.\n\n"
                "## LIMITATIONS & COMPLIANCE POLICIES\n{policy_block}\n  - No confidential or PII data\n  - Follow {industry} industry standards\n\n"
                "## OUTPUT FORMAT\n1. Executive Summary\n2. Main Content\n3. Key Recommendations\n4. Compliance Notes",
                "Initial CORLO template", datetime.utcnow().isoformat(), "system"
            )
        )
    conn.commit()
    conn.close()


# ══════════════════════════════════════════════════════════════════════════════
# AZURE CALL HELPER
# ══════════════════════════════════════════════════════════════════════════════
def _azure_chat(messages: list, max_tokens: int = 512, temperature: float = 0.0) -> tuple:
    """Calls Azure OpenAI and returns (content_text, total_tokens). Raises on failure."""
    resp = _azure_client.chat.completions.create(
        model=_AZURE_DEPLOYMENT,
        messages=messages,
        max_tokens=max_tokens,
        temperature=temperature,
    )
    content = resp.choices[0].message.content or ""
    tokens  = resp.usage.total_tokens if resp.usage else 0
    return content, tokens


def call_llm(system_prompt: str, user_prompt: str, max_tokens: int = 1024, temperature: float = 0.4) -> str:
    """
    Public helper for direct LLM calls (e.g. the /api/refine endpoint in routes.py).
    Returns the response text as a string.
    Raises RuntimeError if Azure is not configured and no fallback is possible.
    """
    if HAS_AZURE and _azure_client:
        content, _ = _azure_chat(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user",   "content": user_prompt},
            ],
            max_tokens=max_tokens,
            temperature=temperature,
        )
        return content
    else:
        # Graceful fallback when Azure env vars are not set
        return (
            "[Demo Mode — Azure OpenAI not configured]\n\n"
            "Refinement requires AZURE_OPENAI_API_KEY, AZURE_OPENAI_BASE_URL, and "
            "AZURE_OPENAI_DEPLOYMENT environment variables to be set.\n\n"
            f"Your comment was received: \"{user_prompt[:200]}...\""
        )


def call_llm_messages(messages: list, max_tokens: int = 1024, temperature: float = 0.3) -> str:
    """
    Send a full OpenAI-format message array to Azure OpenAI.
    Used by the Clarifier Agent to preserve proper conversation turns
    (system / user / assistant) so the LLM retains full context.

    Args:
        messages:   List of {"role": ..., "content": ...} dicts.
        max_tokens: Upper bound on generated tokens.
        temperature: Sampling temperature.

    Returns:
        The model's text reply as a string.
    """
    if HAS_AZURE and _azure_client:
        content, _ = _azure_chat(
            messages=messages,
            max_tokens=max_tokens,
            temperature=temperature,
        )
        return content
    else:
        return "[Demo Mode — Azure OpenAI not configured]"


# ══════════════════════════════════════════════════════════════════════════════
# NODE 1 — INTENT + INDUSTRY CLASSIFIER
# ══════════════════════════════════════════════════════════════════════════════
VALID_INTENTS = [
    "proposal", "report", "email", "code", "content",
    "data analysis", "legal", "it support", "hr", "general"
]

VALID_INDUSTRIES = [
    "banking", "healthcare", "retail", "technology", "manufacturing", "general"
]

_INTENT_KEYWORDS = {
    "proposal":      ["proposal", "pitch", "offer", "bid", "rfp", "quotation"],
    "report":        ["report", "summary", "analysis", "findings", "review"],
    "email":         ["email", "mail", "message", "reply", "respond"],
    "code":          ["code", "script", "function", "program", "debug", "fix", "build", "refactor", "test"],
    "content":       ["blog", "article", "post", "content", "write", "draft", "copy"],
    "data analysis": ["analyze", "data", "insights", "chart", "trend", "metric", "dashboard", "kpi"],
    "legal":         ["contract", "legal", "compliance", "agreement", "terms", "policy", "clause"],
    "it support":    ["ticket", "incident", "issue", "support", "itsm", "helpdesk", "outage"],
    "hr":            ["hr", "employee", "leave", "payroll", "onboarding", "performance"],
}

_INDUSTRY_KEYWORDS = {
    "banking":       ["bank", "financial", "finance", "loan", "credit", "investment", "treasury"],
    "healthcare":    ["health", "medical", "hospital", "patient", "clinical", "pharma"],
    "retail":        ["retail", "store", "customer", "ecommerce", "product", "inventory"],
    "technology":    ["tech", "software", "it", "digital", "api", "system", "cloud"],
    "manufacturing": ["manufacturing", "production", "supply chain", "procurement", "warehouse"],
}


def _keyword_fallback_classify(text: str, task_type: str = None):
    text_lower = text.lower()

    # task_type from UI maps directly to intents — use it as first priority
    TASK_TYPE_TO_INTENT = {
        "research":      "report",
        "writing":       "content",
        "strategy":      "proposal",
        "data":          "data analysis",
        "code":          "code",
        "creative":      "content",
        "communication": "email",
        "learning":      "general",
        "automate":      "code",
        "decision":      "report",
    }
    detected_intent = TASK_TYPE_TO_INTENT.get(task_type, None)

    # If task_type didn't give us an intent, fall back to keyword scan
    if not detected_intent:
        detected_intent = "general"
        for intent, keywords in _INTENT_KEYWORDS.items():
            if any(kw in text_lower for kw in keywords):
                detected_intent = intent
                break

    detected_industry = "general"
    for industry, keywords in _INDUSTRY_KEYWORDS.items():
        if any(kw in text_lower for kw in keywords):
            detected_industry = industry
            break

    return detected_intent, detected_industry


def classify_intent(state: OrchestratorState) -> OrchestratorState:
    role        = state.get("role", "general")
    task_type   = state.get("task_type", "general")
    sensitivity = state.get("data_sensitivity", "general")

    if HAS_AZURE and _azure_client:
        try:
            classifier_prompt = f"""You are an enterprise task classifier. Analyze the user's request and return ONLY a JSON object.

USER REQUEST: "{state['user_input']}"

USER CONTEXT (use this to sharpen your classification):
- Role: {role}  (e.g. a developer asking something is likely a 'code' intent; an executive asking something is likely a 'report' or 'proposal')
- Task Type selected by user: {task_type}  (treat this as a strong hint for intent)
- Data Sensitivity: {sensitivity}  (client = confidential work; internal = inside the org; general = public)

CLASSIFY into exactly one INTENT from this list:
- proposal     : creating proposals, pitches, bids, RFPs, client offers, deliverables, presentations for clients
- report       : reports, summaries, analysis documents, findings, executive reviews, status updates
- email        : writing emails, messages, replies, follow-ups, communications
- code         : programming, scripting, debugging, refactoring, testing, DevOps, APIs, automation
- content      : blog posts, articles, marketing copy, social media, product descriptions, creative writing
- data analysis: data insights, dashboards, KPIs, metrics, charts, business intelligence, trend analysis
- legal        : contracts, compliance, legal review, agreements, terms, regulatory, policy documents
- it support   : IT tickets, incidents, outages, helpdesk, ITSM, change requests, infrastructure issues
- hr           : HR tasks, employee management, payroll, onboarding, performance reviews, leave requests
- general      : anything that does not clearly fit the above categories

CLASSIFY into exactly one INDUSTRY from this list:
- banking      : banking, finance, fintech, insurance, investment, wealth management, treasury, capital markets
- healthcare   : healthcare, medical, hospital, pharma, clinical, biotech, patient care, health IT
- retail       : retail, e-commerce, consumer goods, supply chain, merchandise, stores, omnichannel
- technology   : software, tech, IT services, SaaS, cloud, cybersecurity, data engineering, platforms
- manufacturing: manufacturing, production, industrial, automotive, logistics, factory, operations
- general      : cannot determine industry or does not fit above categories

RULES:
- The user's Role and Task Type are strong signals — weight them heavily alongside the request text.
- Read the FULL meaning, not just keywords. "Goldman Sachs integration" = banking. "Patient portal" = healthcare.
- If user mentions company names, infer the industry (SAP = technology, NHS = healthcare, etc.)
- Return ONLY valid JSON. No explanation. No markdown.

JSON FORMAT:
{{
  "intent": "<one of the 10 intents above>",
  "industry": "<one of the 6 industries above>",
  "intent_confidence": "HIGH or MEDIUM or LOW",
  "industry_confidence": "HIGH or MEDIUM or LOW",
  "reasoning": "<one sentence explaining your classification>"
}}"""

            raw, _ = _azure_chat(
                messages=[
                    {"role": "system", "content": "You are a JSON-only enterprise task classifier. Output valid JSON only. No markdown, no preamble, no explanation outside the JSON."},
                    {"role": "user", "content": classifier_prompt}
                ],
                max_tokens=200,
                temperature=0.0,
            )
            raw    = raw.replace("```json", "").replace("```", "").strip()
            data   = json.loads(raw)
            intent   = data.get("intent", "general")
            industry = data.get("industry", "general")
            if intent not in VALID_INTENTS:
                intent = "general"
            if industry not in VALID_INDUSTRIES:
                industry = "general"
            return {**state, "intent": intent, "industry": industry}

        except Exception:
            intent, industry = _keyword_fallback_classify(state["user_input"], state.get("task_type"))
            return {**state, "intent": intent, "industry": industry}
    else:
        intent, industry = _keyword_fallback_classify(state["user_input"], state.get("task_type"))
        return {**state, "intent": intent, "industry": industry}


# ══════════════════════════════════════════════════════════════════════════════
# NODE 2 — AI TOOL RECOMMENDER
# ══════════════════════════════════════════════════════════════════════════════
def _role_matches(user_role: str, tool_roles: list) -> bool:
    """
    Return True if the user's role appears in the tool's allowed roles list.
    An empty roles list means the tool is open to ALL roles.
    """
    if not tool_roles:
        return True
    u = user_role.lower()
    return any(u in r.lower() or r.lower() in u for r in tool_roles)


def _score_tools_from_registry(user_input: str, intent: str, role: str) -> str:
    """
    Pure registry scoring — used ONLY as a last-resort fallback when the LLM
    call itself fails (e.g. Azure is down, JSON parse error).
    Uses enriched summaries if available, otherwise Excel description.

    This is NOT the primary recommendation path. The primary path is the LLM
    reasoning over enriched tool profiles in recommend_tool().
    """
    if not AI_TOOLS_REGISTRY:
        return ""

    eligible     = _filter_eligible_tools(user_input, intent)
    text_lower   = user_input.lower()
    intent_lower = intent.lower()
    scores       = {}

    for name, info in AI_TOOLS_REGISTRY.items():
        if name not in eligible:
            continue
        score = 0

        # Role match — primary signal (+3)
        if _role_matches(role, info.get("roles", [])):
            score += 3

        # Strong signals — highest-precision keyword match (+3 each)
        for kw in info.get("strong_signals", []):
            if kw.lower() in text_lower:
                score += 3

        # Best-for phrases — broader match (+2 each)
        for kw in info.get("best_for", []):
            if kw.lower() in text_lower:
                score += 2

        # Intent match against enriched summary + category + description (+1)
        enriched  = _get_enriched_summary(name).lower()
        searchable = (
            info.get("category", "") + " " +
            info.get("description", "") + " " +
            enriched
        ).lower()
        if intent_lower in searchable:
            score += 1

        # Disqualify if user input matches a not_for phrase (-99)
        for nf in info.get("not_for", []):
            kw = nf.lower()
            if kw and (kw in text_lower or kw in intent_lower):
                score -= 99

        scores[name] = score

    best_name  = max(scores, key=scores.get)
    best_score = scores[best_name]

    if best_score <= 0:
        role_matched = [
            n for n, info in AI_TOOLS_REGISTRY.items()
            if _role_matches(role, info.get("roles", []))
        ]
        if role_matched:
            return role_matched[0]
        return next(iter(AI_TOOLS_REGISTRY))

    return best_name


def _filter_eligible_tools(user_input: str, intent: str) -> set:
    """
    Pre-filter: return only tools whose output_type is plausibly what the user
    needs, based on keyword matching against output_type_keywords.

    Logic:
      1. For each tool, check if ANY of its output_type_keywords appear in the
         combined user_input + intent text.
      2. If at least one keyword matches → tool is eligible.
      3. If a tool has NO output_type_keywords defined (legacy row) → always eligible
         (fail-open so nothing is silently excluded).
      4. If NO tools match at all (very unusual input) → return all tools (fail-open).

    This ensures tools like GenAI Amplifier (output: wricef_document) are only
    included when the user explicitly asks for wricef / delivery artifacts / teams
    transcript conversion — not for generic proposal or RFP tasks.
    """
    text = (user_input + " " + intent).lower()
    eligible: set = set()

    for name, info in AI_TOOLS_REGISTRY.items():
        kw_list = info.get("output_type_keywords", [])
        if not kw_list:
            eligible.add(name)
            continue
        for kw in kw_list:
            if kw.lower() in text:
                eligible.add(name)
                break

    if not eligible:
        return set(AI_TOOLS_REGISTRY.keys())

    return eligible


# ══════════════════════════════════════════════════════════════════════════════
# AGENT 3 — INTERNAL TOOL RECOMMENDER
#
# Exact prompt from Agentic AI Navigator.
# Queries tool_knowledge ChromaDB (vector search) for the top-K most relevant
# internal tool chunks, then asks the LLM to score and rank top 3.
# Only operates on tools where is_internal=True.
# ══════════════════════════════════════════════════════════════════════════════

def _build_internal_agent_prompt(relevant_tools_json: str) -> str:
    return f"""
You are the Internal Tool Recommender Agent.

Your job is to deeply understand the user's task and recommend the top 3 most relevant tools exclusively from the INTERNAL TOOL CATALOG provided below.

-------------------------
STEP 1 — UNDERSTAND THE TASK:
Before recommending anything, carefully analyze the user's task:
- Identify the user's role (e.g., Developer, Marketer, HR, Analyst)
- Identify the core intent and goal of the task (what outcome they want)
- Identify the specific parameters (e.g., scale, format, domain, workflow type)
- Understand the full context — do not make surface-level matches

-------------------------
STEP 2 — UNDERSTAND THE INTERNAL TOOL CATALOG:
Read and understand every tool in the INTERNAL TOOL CATALOG below thoroughly:
- Understand each tool's name, capabilities, use cases, supported roles, and features
- Do not skim — each tool's details matter for accurate matching
- NOTE: This catalog contains only the most semantically relevant tools for the user's task, pre-filtered using vector similarity search.

-------------------------
STEP 3 — MATCH & RECOMMEND:
Based on your understanding of both the task and the catalog:
- Select the TOP 3 tools that best fit the user's task, intent, and role
- Rank them from most relevant to least relevant
- Assign a relevance score from 0 to 100 for each tool based on how well it fits
- Provide a specific, clear reason for each recommendation that ties back to the user's actual task

-------------------------
CRITICAL RULES:
1. You MUST ONLY recommend tools that exist in the INTERNAL TOOL CATALOG below. Absolutely no hallucinations or external tools.
2. Never recommend a tool just because its name sounds relevant — validate against its actual catalog description and capabilities.
3. Always return exactly 3 tools with Tool Name, Score, and Reason. If fewer than 3 tools are relevant, clearly note it but still return the best available.
4. Scores must reflect true relevance — avoid giving identical scores unless genuinely equal.
5. Reasons must be specific to the user's role and task — not generic descriptions of the tool.

-------------------------
INTERNAL TOOL CATALOG (Top-5 Semantically Relevant Tools Retrieved via Vector Search):
{relevant_tools_json}

Return ONLY this JSON (no markdown, no extra text):
{{
  "recommendations": [
    {{"tool": "<tool name>", "score": <0-100>, "reason": "<specific justification>"}},
    {{"tool": "<tool name>", "score": <0-100>, "reason": "<specific justification>"}},
    {{"tool": "<tool name>", "score": <0-100>, "reason": "<specific justification>"}}
  ]
}}
"""


def internal_agent(state: OrchestratorState) -> OrchestratorState:
    """
    Agent 3 — Internal Tool Recommender.
    Filters registry to internal tools only, queries ChromaDB for semantic
    matches, then asks the LLM to pick top 3 with scores and reasons.
    """
    user_input = state["user_input"]
    user_role  = state.get("role", "general").strip()
    intent     = state.get("intent", "general").strip()

    internal_tools = {
        name: info for name, info in AI_TOOLS_REGISTRY.items()
        if info.get("is_internal", False)
    }

    if not internal_tools:
        return {**state, "internal_results": json.dumps({"recommendations": []})}

    eligible = _filter_eligible_tools(user_input, intent)
    eligible_internal = {
        name: info for name, info in internal_tools.items()
        if name in eligible and _role_matches(user_role, info.get("roles", []))
    } or internal_tools

    tool_names = list(eligible_internal.keys())

    knowledge = {}
    try:
        knowledge = query_tool_knowledge(user_input, tool_names, n_results=5)
    except Exception:
        pass

    catalog_entries = []
    for name, info in eligible_internal.items():
        entry = {
            "tool_name":   name,
            "description": info.get("description", ""),
            "category":    info.get("category", ""),
            "best_for":    info.get("best_for", []),
            "strong_signals": info.get("strong_signals", []),
            "not_for":     info.get("not_for", []),
            "knowledge_excerpts": knowledge.get(name, [])[:2],
        }
        catalog_entries.append(entry)

    relevant_tools_json = json.dumps(catalog_entries, indent=2)
    system_prompt = _build_internal_agent_prompt(relevant_tools_json)

    if HAS_AZURE and _azure_client:
        try:
            raw, _ = _azure_chat(
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user",   "content": f"Task:\n{user_input}\nRole: {user_role}"},
                ],
                max_tokens=800,
                temperature=0.1,
            )
            raw = (raw or "").replace("```json", "").replace("```", "").strip()
            json.loads(raw)
            return {"internal_results": raw}
        except Exception as e:
            fallback = _score_tools_from_registry(user_input, intent, user_role)
            fallback_result = json.dumps({
                "recommendations": [
                    {"tool": fallback, "score": 70, "reason": f"Selected via registry scoring (LLM error: {str(e)[:60]})."}
                ]
            })
            return {"internal_results": fallback_result}

    fallback = _score_tools_from_registry(user_input, intent, user_role)
    return {"internal_results": json.dumps({
        "recommendations": [
            {"tool": fallback, "score": 60, "reason": "Selected via keyword registry scoring (no LLM configured)."}
        ]
    })}


# ══════════════════════════════════════════════════════════════════════════════
# AGENT 4 — EXTERNAL TOOL RECOMMENDER
#
# Exact prompt from Agentic AI Navigator.
# Operates on tools where is_internal=False — uses LLM's own knowledge.
# No vector search needed — passes tool names + profiles directly.
# ══════════════════════════════════════════════════════════════════════════════

def _build_external_agent_prompt(tool_names: list) -> str:
    if not tool_names:
        numbered = "(no external tools configured)"
    else:
        numbered = "\n".join(f"{i+1}. {name}" for i, name in enumerate(tool_names))

    return f"""
You are an intelligent External Tool Recommendation Agent.

Your job is to analyze a user's query or task, understand the user's role and the context of their task, and recommend the top 3 most suitable tools from a predefined list.

-------------------------
AVAILABLE TOOLS:
You MUST ONLY recommend tools from the following list:

{numbered}

Do NOT recommend any tool outside this list under any circumstances.

-------------------------
YOUR OBJECTIVE:

1. UNDERSTAND THE USER:
   - Identify the user's role (explicit or inferred)
   - Identify the core task or goal
   - Identify task type (e.g., content creation, CRM, automation, legal contracts, AI assistance, etc.)

2. USE YOUR KNOWLEDGE:
   - Use your internal knowledge about the strengths and use-cases of each tool
   - Match tools based on how well they solve the user's problem
   - Consider efficiency, relevance, and specialization

3. SELECT TOP 3 TOOLS:
   - Choose exactly 3 tools from the list
   - Rank them from most suitable to least suitable

4. SCORING:
   - Assign a score to each tool between 0 and 100
   - Scores must reflect relevance and effectiveness for the task
   - The best tool should have the highest score
   - Avoid giving identical scores unless absolutely necessary

5. JUSTIFICATION:
   - Provide a short, clear reason why each tool is selected
   - Justifications must be specific to the user's role and task

OUTPUT FORMAT (STRICT JSON — no markdown, no extra text):
{{
  "user_role": "<identified role>",
  "task_summary": "<brief summary of the task>",
  "recommendations": [
    {{"tool": "<tool name>", "score": <number>, "reason": "<concise justification>"}},
    {{"tool": "<tool name>", "score": <number>, "reason": "<concise justification>"}},
    {{"tool": "<tool name>", "score": <number>, "reason": "<concise justification>"}}
  ]
}}

RULES:
- Only output valid JSON
- Always return exactly 3 tools
- Never hallucinate tools outside the provided list
- If user role is not explicitly stated, infer it from the task
"""


def _load_external_tools_json() -> list:
    """
    Load external tools from data/external_tools.json.
    Supports both formats:
      - Old: ["ChatGPT", ...]
      - New: [{"name": "ChatGPT", "url": "https://..."}, ...]
    Returns list of {"name": ..., "url": ...} dicts.
    Falls back to registry tools where is_internal=False if file is missing.
    """
    json_path = os.path.join(os.path.dirname(__file__), "data", "external_tools.json")
    try:
        with open(json_path, "r") as f:
            data = json.load(f)
        tools = data.get("tools", [])
        if tools:
            normalised = []
            for t in tools:
                if isinstance(t, dict):
                    normalised.append({"name": t.get("name", ""), "url": t.get("url", "")})
                else:
                    normalised.append({"name": str(t), "url": ""})
            return [t for t in normalised if t["name"]]
    except Exception:
        pass
    return [
        {"name": name, "url": info.get("url", "")}
        for name, info in AI_TOOLS_REGISTRY.items()
        if not info.get("is_internal", False)
    ]


def external_agent(state: OrchestratorState) -> OrchestratorState:
    """
    Agent 4 — External Tool Recommender.
    Loads tool names from data/external_tools.json (same as Agentic Navigator).
    Falls back to registry tools with is_internal=False if file not found.
    Passes tool names to the LLM which uses its own knowledge to rank top 3.
    """
    user_input  = state["user_input"]
    user_role   = state.get("role", "general").strip()

    tools       = _load_external_tools_json()
    url_map     = {t["name"]: t["url"] for t in tools}
    names_only  = [t["name"] for t in tools]

    if not names_only:
        return {"external_results": json.dumps({"recommendations": [], "url_map": {}})}

    system_prompt = _build_external_agent_prompt(names_only)

    def _inject_urls(raw_json: str) -> str:
        try:
            data = json.loads(raw_json)
            for rec in data.get("recommendations", []):
                rec["url"] = url_map.get(rec.get("tool", ""), "")
            data["url_map"] = url_map
            return json.dumps(data)
        except Exception:
            return raw_json

    if HAS_AZURE and _azure_client:
        try:
            raw, _ = _azure_chat(
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user",   "content": f"Task:\n{user_input}\nRole: {user_role}"},
                ],
                max_tokens=800,
                temperature=0.1,
            )
            raw = (raw or "").replace("```json", "").replace("```", "").strip()
            return {"external_results": _inject_urls(raw)}
        except Exception as e:
            first = names_only[0] if names_only else "ChatGPT"
            return {"external_results": _inject_urls(json.dumps({
                "recommendations": [
                    {"tool": first, "score": 65,
                     "reason": f"Selected as best available external tool (LLM error: {str(e)[:60]})."}
                ]
            }))}

    first = names_only[0] if names_only else "ChatGPT"
    return {"external_results": _inject_urls(json.dumps({
        "recommendations": [
            {"tool": first, "score": 55, "reason": "Selected via registry (no LLM configured)."}
        ]
    }))}


# ══════════════════════════════════════════════════════════════════════════════
# AGENT 5 — FINAL DECIDER
#
# Exact prompt from Agentic AI Navigator.
# Merges top-3 from Agent 3 (internal) + top-3 from Agent 4 (external) = 6 tools.
# Picks the overall best top 3 and returns clean Markdown.
# Also populates recommended_tool / tool_alternatives for backward compatibility
# with the rest of your pipeline (audit log, CORLO prompt, etc).
# ══════════════════════════════════════════════════════════════════════════════

_FINAL_DECIDER_PROMPT = """
You are the Final Decider Agent.
You will be provided with:
1. The user's original detailed task.
2. The Top 3 recommendations (with scores) from the Internal Agent.
3. The Top 3 recommendations (with scores) from the External Agent.

CRITICAL RULES:
1. Carefully read and understand the user's task, role, and requirements in full before making any decision.
2. Compare all 6 recommended tools holistically — consider both the scores AND the reasons provided by each agent to judge true relevance to the user's task.
3. Select the absolute TOP 3 tools overall. Sort them strictly by score from highest to lowest (the highest scoring tool must appear first).
4. Do NOT mention whether a tool is internal or external. Do not include labels like [Internal Tool] or [External Tool] or any similar badge/tag.
5. For each tool, display:
   - The tool name (bold)
   - Its score (e.g., Score: 92/100)
   - A clear, professional justification that is specific to the user's role and task
6. Do not include any other text, preamble, or commentary outside the formatted top 3 list.
7. Also return a JSON block at the end (after the markdown) in this exact format so the system can parse it:
---JSON---
{"top1": {"tool": "<name>", "score": <int>, "reason": "<reason>"}, "top2": {"tool": "<name>", "score": <int>, "reason": "<reason>"}, "top3": {"tool": "<name>", "score": <int>, "reason": "<reason>"}}

OUTPUT FORMAT (Markdown first, then JSON block):
## Final Top 3 Recommended AI Tools

**1. <Tool Name>** — Score: <score>/100
<Justification specific to the user's task>

**2. <Tool Name>** — Score: <score>/100
<Justification specific to the user's task>

**3. <Tool Name>** — Score: <score>/100
<Justification specific to the user's task>
---JSON---
{"top1": ..., "top2": ..., "top3": ...}
"""


def final_decider(state: OrchestratorState) -> OrchestratorState:
    """
    Agent 5 — Final Decider.
    Merges internal + external results, picks overall top 3, returns Markdown.
    Also sets recommended_tool / tool_alternatives on state for the rest of
    the pipeline (audit log, CORLO prompt builder, etc.).
    """
    internal_raw = state.get("internal_results") or "{}"
    external_raw = state.get("external_results") or "{}"

    # Build url_map from external results so Open buttons work
    url_map = {}
    try:
        url_map = json.loads(external_raw).get("url_map", {})
    except Exception:
        pass

    combined_input = (
        f"USER TASK:\n{state['user_input']}\n\n"
        f"--- INTERNAL AGENT RECOMMENDATIONS ---\n{internal_raw}\n\n"
        f"--- EXTERNAL AGENT RECOMMENDATIONS ---\n{external_raw}"
    )

    def _parse_fallback() -> dict:
        top_tool = ""
        top_reason = ""
        alts = []
        alt_reasons = []
        alt_pcts = []
        alt_urls = []
        for raw in [internal_raw, external_raw]:
            try:
                data = json.loads(raw)
                recs = data.get("recommendations", [])
                for r in recs:
                    tool_url = url_map.get(r.get("tool", ""), r.get("url", ""))
                    if not top_tool:
                        top_tool   = r.get("tool", "")
                        top_reason = r.get("reason", "")
                    else:
                        alts.append(r.get("tool", ""))
                        alt_reasons.append(r.get("reason", ""))
                        alt_pcts.append(int(r.get("score", 50)))
                        alt_urls.append(tool_url)
            except Exception:
                pass
        return {
            "recommended_tool":                 top_tool,
            "tool_url":                         url_map.get(top_tool, ""),
            "tool_reason":                      top_reason,
            "tool_confidence":                  "MEDIUM",
            "tool_confidence_pct":              50,
            "tool_confidence_explanation":      top_reason,
            "tool_alternatives":                alts[:2],
            "tool_alternative_reasons":         alt_reasons[:2],
            "tool_alternative_confidence_pcts": alt_pcts[:2],
            "tool_alternative_urls":            alt_urls[:2],
        }

    if HAS_AZURE and _azure_client:
        try:
            raw, _ = _azure_chat(
                messages=[
                    {"role": "system", "content": _FINAL_DECIDER_PROMPT},
                    {"role": "user",   "content": combined_input},
                ],
                max_tokens=1000,
                temperature=0.1,
            )
            raw = (raw or "").strip()

            markdown_part = raw
            parsed_tools = {}
            if "---JSON---" in raw:
                parts = raw.split("---JSON---", 1)
                markdown_part = parts[0].strip()
                try:
                    parsed_tools = json.loads(parts[1].strip())
                except Exception:
                    pass

            top1 = parsed_tools.get("top1", {})
            top2 = parsed_tools.get("top2", {})
            top3 = parsed_tools.get("top3", {})

            recommended_tool = top1.get("tool", "")
            tool_reason      = top1.get("reason", "")
            tool_score       = int(top1.get("score", 80))
            confidence       = "HIGH" if tool_score >= 75 else "MEDIUM" if tool_score >= 50 else "LOW"

            alts        = [t.get("tool", "") for t in [top2, top3] if t.get("tool")]
            alt_reasons = [t.get("reason", "") for t in [top2, top3] if t.get("tool")]
            alt_pcts    = [int(t.get("score", 50)) for t in [top2, top3] if t.get("tool")]
            alt_urls    = [url_map.get(t.get("tool", ""), "") for t in [top2, top3] if t.get("tool")]

            if not recommended_tool:
                fallback = _parse_fallback()
                return {**fallback, "llm_output": markdown_part}

            return {
                "recommended_tool":                 recommended_tool,
                "tool_url":                         url_map.get(recommended_tool, ""),
                "tool_reason":                      tool_reason,
                "tool_confidence":                  confidence,
                "tool_confidence_pct":              tool_score,
                "tool_confidence_explanation":      tool_reason,
                "tool_alternatives":                alts,
                "tool_alternative_reasons":         alt_reasons,
                "tool_alternative_confidence_pcts": alt_pcts,
                "tool_alternative_urls":            alt_urls,
            }
        except Exception:
            pass

    return _parse_fallback()


# ══════════════════════════════════════════════════════════════════════════════
# POLICY CHECKER AGENT — Agent 2
#
# Exact prompt from Agentic AI Navigator agents_graph.py.
# Retrieves top-K policy chunks via ChromaDB vector search, then asks the LLM
# to output ONLY [POLICY_VIOLATED] + Reason OR [POLICY_CLEAR].
# Replaces the old retrieve_policies + check_policy_compliance nodes.
# ══════════════════════════════════════════════════════════════════════════════

_POLICY_CHECKER_SYSTEM_PROMPT = """
You are a strict Enterprise Policy Compliance Agent.

Your job is to carefully review the user's task against the provided POLICY GUIDELINES and determine whether the task contains or relates to any prohibited topics, restricted words, or guideline violations.

-------------------------
INSTRUCTIONS:

1. Read the user's task carefully.
2. Read every POLICY GUIDELINE provided below.
3. Determine if the task touches upon, implies, or relates to ANY prohibited topic or restricted content.
4. Output ONLY one of the following two exact formats — no extra text:

If a violation IS found:
[POLICY_VIOLATED]
Reason: <short explanation of exactly which policy guideline is violated and why>

If NO violation is found:
[POLICY_CLEAR]

-------------------------
CRITICAL RULES:
- Be strict and thorough — err on the side of caution.
- Even indirect or implied references to prohibited topics count as violations.
- Do not add any commentary, explanation, or text outside the two formats above.
- Do not suggest alternatives or ask questions.
"""

_POLICY_TOP_K = 5


def retrieve_policies(state: OrchestratorState) -> OrchestratorState:
    """
    NODE 3 — Policy RAG Retrieval.
    Queries ChromaDB with the user task to find the top-K most relevant
    policy chunks. Stores them in state['policies'] for Agent 2.
    """
    query = f"{state['intent']} {state['industry']} {state['user_input']}"
    try:
        count = policy_collection.count()
        if count == 0:
            return {**state, "policies": []}
        effective_k = min(_POLICY_TOP_K, count)
        results  = policy_collection.query(query_texts=[query], n_results=effective_k)
        docs     = results.get("documents", [[]])[0]
        policies = docs if docs else []
    except Exception:
        policies = []
    return {**state, "policies": policies}


def check_policy_compliance(state: OrchestratorState) -> OrchestratorState:
    """
    NODE 3b — Policy Checker Agent (Agent 2).
    Feeds the clarified task + retrieved policy chunks to the LLM using the
    strict Agent 2 prompt. Parses [POLICY_VIOLATED] / [POLICY_CLEAR] output.
    Sets policy_blocked=True when a violation is found.
    """
    policies       = state.get("policies", [])
    user_input     = state["user_input"]
    existing_flags = state.get("policy_flags", [])

    if not policies:
        return {
            **state,
            "policy_summary": (
                "No company policy documents have been uploaded yet. "
                "General enterprise best practices apply."
            ),
            "policy_blocked": False,
            "policy_flags":   existing_flags,
        }

    policy_context = "\n".join(f"- {p}" for p in policies)

    user_msg = (
        f"USER TASK:\n{user_input}\n\n"
        f"POLICY GUIDELINES (most relevant to this task):\n{policy_context}"
    )

    if HAS_AZURE and _azure_client:
        try:
            raw, _ = _azure_chat(
                messages=[
                    {"role": "system", "content": _POLICY_CHECKER_SYSTEM_PROMPT},
                    {"role": "user",   "content": user_msg},
                ],
                max_tokens=300,
                temperature=0.0,
            )
            raw = (raw or "").strip()

            if "[POLICY_VIOLATED]" in raw:
                reason = (
                    raw.replace("[POLICY_VIOLATED]", "")
                       .replace("Reason:", "")
                       .strip()
                )
                return {
                    **state,
                    "policy_summary": reason,
                    "policy_blocked": True,
                    "policy_flags":   list(set(existing_flags + [f"Policy violation: {reason[:120]}"])),
                }

            return {
                **state,
                "policy_summary": "Task reviewed — no policy violations found.",
                "policy_blocked": False,
                "policy_flags":   existing_flags,
            }

        except Exception:
            pass

    return {
        **state,
        "policy_summary": (
            "Policy check could not be completed (LLM unavailable). "
            "Proceeding with general enterprise best practices."
        ),
        "policy_blocked": False,
        "policy_flags":   existing_flags,
    }
















# ══════════════════════════════════════════════════════════════════════════════
def _build_system_prompt(role: str, task_type: str, sensitivity: str,
                          industry: str, intent: str, tool_name: str,
                          tool_info: dict) -> str:
    """
    Builds the SYSTEM prompt that tells the LLM WHO it is and HOW to behave.
    The role from the user (matched against Excel) drives the persona and tone.
    """
    effective_role = role.strip() if role and role != "general" else "expert enterprise professional"

    # Per-role behavioural instructions — derived from what each role needs
    role_behaviour = {
        "Executive / Director": (
            "Present insights at a strategic level. Be concise and outcome-focused. "
            "Lead with key decisions, business impact, and ROI. Avoid deep technical detail. "
            "Use executive-friendly language: bullet summaries, clear headers, no jargon."
        ),
        "Business Analyst": (
            "Provide structured, data-backed analysis. Use tables, bullet points, and numbered lists. "
            "Highlight assumptions, gaps, and recommendations clearly. "
            "Balance business context with analytical rigour."
        ),
        "Developer / Technical": (
            "Be technically precise and detailed. Include code snippets, commands, configurations, "
            "or architecture diagrams where relevant. Assume strong technical literacy. "
            "Use correct technical terminology. Format code in proper blocks."
        ),
        "Consultant / Manager": (
            "Balance technical accuracy with business clarity. Highlight risks, timelines, "
            "dependencies, and stakeholder considerations. Structure output for client-ready delivery. "
            "Use professional consulting language."
        ),
        "Finance / Accounting": (
            "Prioritise numerical accuracy and compliance. Use structured tabular formats where possible. "
            "Flag any figures that require validation. Align with accounting standards. "
            "Avoid ambiguous language around financial figures."
        ),
        "HR / People Ops": (
            "Use empathetic, people-first language. Ensure tone is inclusive and policy-compliant. "
            "Avoid jargon. Structure content to be accessible to all employee levels. "
            "Highlight any legal or HR compliance considerations."
        ),
        "Sales / BD": (
            "Emphasise value propositions, client benefits, and competitive differentiators. "
            "Keep tone persuasive, confident, and professional. "
            "Focus on outcomes, ROI, and solving client pain points. "
            "Structure for use in client-facing communications."
        ),
        "Marketing / Comms": (
            "Prioritise clarity, brand voice, and audience engagement. "
            "Structure content for readability and impact. "
            "Use compelling language appropriate for the target audience. "
            "Adapt tone based on channel (internal vs. external, formal vs. casual)."
        ),
    }.get(effective_role,
          "Provide a clear, professional, well-structured response appropriate for the user's context.")

    # Sensitivity-based content rules
    sensitivity_rules = {
        "client": (
            "⚠️ CONFIDENTIAL / CLIENT DATA RULES:\n"
            "- Replace all real names with [CLIENT NAME], [CONTACT NAME]\n"
            "- Replace specific figures with [VALUE] or [AMOUNT]\n"
            "- Do NOT reproduce any PII, account numbers, or contract specifics\n"
            "- Flag any section that requires human review before sharing externally"
        ),
        "internal": (
            "🔒 INTERNAL DATA RULES:\n"
            "- Use general terms for sensitive internal metrics\n"
            "- Do not disclose specific internal figures that could be sensitive if leaked\n"
            "- Mark any section intended for internal use only"
        ),
        "general": (
            "✅ GENERAL DATA: Standard professional best practices apply. "
            "No special masking required."
        ),
    }.get(sensitivity, "Standard data handling applies.")

    # Tool-specific usage hint
    tool_hint = (
        f"The output will be used in {tool_name} ({tool_info.get('category', 'AI Tool')}). "
        f"Structure and format the response to be directly usable in that tool."
    )

    return f"""You are a {effective_role} operating in the {industry} industry.

BEHAVIOURAL INSTRUCTIONS FOR THIS ROLE:
{role_behaviour}

TOOL CONTEXT:
{tool_hint}

DATA HANDLING RULES:
{sensitivity_rules}

GENERAL RULES:
- Return ONLY the final response — no meta-commentary, no preamble, no "here is your response"
- Always define clear sections with headers
- Always include a Compliance / Risk note at the end
- Tailor depth and tone exactly to the role described above
- Flag any area that requires human expert review"""


def _build_user_prompt(state: OrchestratorState, tool_info: dict, policy_block: str,
                        prompt_version: str) -> str:
    """
    Builds a dynamic, task-adaptive CORLO prompt.

    Instead of a fixed 5-section template, this function selects the right
    prompt *shape* based on the combination of intent + task_type + role.
    Each shape has different sections, different depth, different tone.

    The goal: a prompt that looks like it was written by an expert specifically
    for this task — not a fill-in-the-blanks form.
    """
    role        = state.get("role",      "general").strip()
    task_type   = state.get("task_type", "general").strip()
    industry    = state["industry"]
    intent      = state["intent"]
    tool_name   = state["recommended_tool"]
    user_input  = state["user_input"]
    effective_role = role if role and role != "general" else "Enterprise Professional"

    # ── Policy block ──────────────────────────────────────────────────────────
    has_policies = bool(
    policy_block and policy_block.strip()
    and "No specific policies" not in policy_block
    and "Policy retrieval unavailable" not in policy_block
    )
    policy_section = (
        f"The following company policies apply. Do not produce output that conflicts "
        f"with them:\n\n{policy_block}"
        if has_policies
        else f"No specific company policies were retrieved. Apply {industry} industry best practices."
    )

    # ═══════════════════════════════════════════════════════════════════════
    # PROMPT SHAPE SELECTOR
    # Maps (intent, task_type) combinations to a prompt architecture.
    # Each shape builds radically different sections.
    # ═══════════════════════════════════════════════════════════════════════

    # ── SHAPE: CODE / TECHNICAL ───────────────────────────────────────────
    if intent == "code" or task_type == "code":
        return f"""You are a senior {effective_role} and software engineer working in the {industry} industry.
You write clean, production-ready code with proper error handling and comments.

**TASK**
{user_input}

**TECHNICAL REQUIREMENTS**
- Language/framework: infer from the request, or ask if ambiguous
- Code quality: production-grade — no pseudocode, no placeholders
- Include: working implementation + inline comments explaining key decisions
- Include: usage example showing how to call/run the code
- Include: any dependencies, prerequisites, or setup steps needed
- Consider: edge cases, error handling, and security implications

**TOOL CONTEXT**
Output is designed for use in {tool_name} ({tool_info.get('category', 'AI Tool')}).

**CONSTRAINTS**
{policy_section}

**OUTPUT STRUCTURE**
1. Brief explanation of your approach (2-3 sentences)
2. Complete, runnable code in properly labelled code blocks
3. Usage example
4. Notes on edge cases, limitations, or things the developer should watch out for
5. Any follow-up steps (e.g. tests to write, configs to set)

Write the response now."""

    # ── SHAPE: EMAIL / COMMUNICATION ─────────────────────────────────────
    if intent == "email" or task_type == "communication":
        return f"""You are a {effective_role} in the {industry} industry drafting a professional communication.
Your output must be ready to copy and send — no placeholders, no rewrites needed.

**TASK**
{user_input}

**COMMUNICATION REQUIREMENTS**
- Tone: professional, appropriate for the relationship and context implied in the request
- Length: as long as the message needs to be — not a word more
- Subject line: include one (infer from context)
- Opening: gets to the point quickly — no filler openers
- Body: clear, well-structured, purpose-driven
- Closing: appropriate call to action or next step

**TOOL CONTEXT**
Optimised for {tool_name}.

**CONSTRAINTS**
{policy_section}

**OUTPUT**
Provide:
- Subject: [subject line]
- [The complete email/message body, ready to send]
- Optional: one-line note on tone choices if the context is nuanced

Write the communication now."""

    # ── SHAPE: DATA ANALYSIS ──────────────────────────────────────────────
    if intent == "data analysis" or task_type == "data":
        return f"""You are a {effective_role} and data analyst working in the {industry} industry.
Your job is to turn data and observations into clear, actionable intelligence.

**ANALYTICAL TASK**
{user_input}

**ANALYSIS REQUIREMENTS**
- Lead with the most important insight — what does this data actually mean?
- Identify patterns, trends, anomalies, and their likely causes
- Quantify where possible; flag where data is missing or assumptions are made
- Separate facts (what the data shows) from interpretation (what it might mean)
- Recommend concrete next steps based on the findings

**TOOL CONTEXT**
Analysis structured for use in {tool_name} ({tool_info.get('category', 'AI Tool')}).

**CONSTRAINTS**
{policy_section}

**OUTPUT STRUCTURE**
1. **Key Finding** — the single most important insight (1-2 sentences)
2. **Analysis** — detailed breakdown with supporting points
3. **Data Gaps / Assumptions** — what's missing, what was assumed
4. **Visualisation Suggestions** — chart types or views that would communicate this well
5. **Recommendations** — 3-5 specific, actionable next steps
6. **Risks / Caveats** — where the analysis could be wrong or misleading

Produce the analysis now."""

    # ── SHAPE: REPORT / SUMMARY ───────────────────────────────────────────
    if intent == "report" or task_type in ("research", "decision"):
        depth_instruction = {
            "research": "Be thorough and evidence-based. Surface insights beyond the obvious.",
            "decision": "Frame as a decision document: options → evaluation → recommendation.",
        }.get(task_type, "Produce a professional report with clear structure and actionable findings.")

        return f"""You are a {effective_role} producing a formal report for the {industry} sector.
{depth_instruction}

**REPORT BRIEF**
{user_input}

**REPORT REQUIREMENTS**
- Audience: {effective_role} and their stakeholders
- Depth: sufficient to inform a decision or communicate findings — not academic
- Structure: clear numbered sections with informative headers
- Evidence: support claims with reasoning; flag where expert validation is needed
- Recommendations: specific and actionable, not generic
- Length: as long as the brief requires — no padding, no omissions

**TOOL CONTEXT**
Formatted for {tool_name}.

**CONSTRAINTS**
{policy_section}

**OUTPUT STRUCTURE**
1. **Executive Summary** — key findings and recommendation in under 150 words
2. **Background / Context** — why this matters
3. **Main Findings** — the substance of the report
4. **Analysis** — what the findings mean
5. **Recommendations** — numbered, specific, owner-assignable
6. **Risks & Mitigations** — what could go wrong
7. **Next Steps** — immediate actions with suggested owners/timelines
8. **Compliance Note** — any areas requiring specialist review

Produce the report now."""

    # ── SHAPE: PROPOSAL / PITCH ───────────────────────────────────────────
    if intent == "proposal" or task_type == "strategy":
        return f"""You are a {effective_role} creating a high-impact proposal for the {industry} industry.
Your output must be persuasive, professional, and immediately presentable.

**PROPOSAL BRIEF**
{user_input}

**PROPOSAL REQUIREMENTS**
- Open with a compelling problem statement — make the reader feel the pain
- Articulate the proposed solution clearly and specifically
- Show the value: quantify benefits where possible (time saved, cost reduced, risk mitigated)
- Address likely objections or concerns pre-emptively
- Close with a clear, confident ask or call to action
- Tone: authoritative, client-ready, outcome-focused

**TOOL CONTEXT**
Built for {tool_name} ({tool_info.get('category', 'AI Tool')}).

**CONSTRAINTS**
{policy_section}

**OUTPUT STRUCTURE**
1. **Executive Summary** — the proposal in 3-4 sentences
2. **Problem Statement** — the challenge being solved
3. **Proposed Solution** — what you're offering and how it works
4. **Value & Benefits** — concrete outcomes and ROI
5. **Approach / Methodology** — how you'll deliver it
6. **Timeline & Milestones** — key phases and dates
7. **Investment / Ask** — what's required (resource, budget, decision)
8. **Why Us / Why Now** — differentiation and urgency
9. **Next Steps** — clear call to action

Produce the proposal now."""

    # ── SHAPE: CONTENT / CREATIVE / MARKETING ────────────────────────────
    if intent == "content" or task_type in ("creative", "writing"):
        audience_hint = f"for a {industry} industry audience" if industry != "general" else "for a professional audience"
        return f"""You are a {effective_role} and content specialist creating original content {audience_hint}.
Your writing is clear, engaging, and purposeful — it earns the reader's attention from the first line.

**CONTENT BRIEF**
{user_input}

**CONTENT REQUIREMENTS**
- Hook: open with something that makes the reader want to continue
- Voice: professional yet human — authoritative without being stiff
- Structure: logical flow, easy to scan, with subheadings where appropriate
- Specificity: concrete details and examples beat vague generalisations
- Length: exactly as long as the content needs to be
- Finish: end with purpose — a takeaway, a question, or a clear next step

**TOOL CONTEXT**
Optimised for {tool_name}.

**CONSTRAINTS**
{policy_section}

**OUTPUT**
Produce the complete, publication-ready content piece.
Do not include meta-commentary about the content — just write it.

Write the content now."""

    # ── SHAPE: LEGAL / COMPLIANCE ─────────────────────────────────────────
    if intent == "legal":
        return f"""You are a {effective_role} with expertise in legal and compliance matters in the {industry} sector.
Note: This output is a professional starting point — it must be reviewed by a qualified legal professional before use.

**LEGAL TASK**
{user_input}

**REQUIREMENTS**
- Accuracy: use correct legal terminology for the {industry} context
- Clarity: make the document understandable to non-lawyers where possible
- Completeness: cover the key clauses/provisions typically needed for this document type
- Flagging: clearly mark any clause that carries significant risk or needs specialist input
- Disclaimers: include appropriate review-required notices

**TOOL CONTEXT**
Drafted using {tool_name}.

**CONSTRAINTS**
{policy_section}

**OUTPUT STRUCTURE**
1. **Document / Clause** — the complete draft
2. **Key Provisions Explained** — plain-English summary of major clauses
3. **Risk Flags** — sections that need legal review highlighted explicitly
4. **⚠️ Legal Disclaimer** — this is a draft only; qualified legal review is required before use

Produce the legal document now."""

    # ── SHAPE: HR ─────────────────────────────────────────────────────────
    if intent == "hr":
        return f"""You are a {effective_role} and HR professional working in the {industry} sector.
You combine deep HR expertise with empathy and clarity, producing people-first documents that are also legally sound.

**HR TASK**
{user_input}

**REQUIREMENTS**
- Tone: inclusive, fair, empathetic — professional but human
- Legal awareness: flag any aspect that may have legal/employment law implications
- Clarity: accessible to all employees, not just HR professionals
- Policy alignment: consistent with standard HR best practices for {industry}

**TOOL CONTEXT**
Formatted for {tool_name}.

**CONSTRAINTS**
{policy_section}

**OUTPUT STRUCTURE**
1. **Main Document / Communication** — the deliverable
2. **Usage Notes** — how and when to use this document
3. **Legal / Compliance Flags** — anything requiring HR or legal review
4. **Recommended Next Steps** — follow-up actions

Produce the HR document now."""

    # ── SHAPE: IT SUPPORT ─────────────────────────────────────────────────
    if intent == "it support":
        return f"""You are a {effective_role} and IT support specialist working in the {industry} sector.
You provide precise, step-by-step technical guidance that resolves issues efficiently.

**IT SUPPORT TASK**
{user_input}

**REQUIREMENTS**
- Diagnose first: identify the most likely root cause(s)
- Be specific: exact commands, settings, or steps — no vague instructions
- Escalation path: when and how to escalate if the primary fix fails
- Prevention: note what can be done to prevent recurrence

**TOOL CONTEXT**
Structured for {tool_name}.

**CONSTRAINTS**
{policy_section}

**OUTPUT STRUCTURE**
1. **Problem Assessment** — likely root cause
2. **Resolution Steps** — numbered, exact steps to resolve
3. **Verification** — how to confirm the fix worked
4. **Escalation Path** — when and how to escalate
5. **Prevention** — steps to avoid recurrence
6. **Impact / Risk Note** — any risks in applying the fix

Produce the IT support response now."""

    # ── SHAPE: GENERAL / FALLBACK — adaptive based on task_type ──────────
    # Even the fallback adapts its framing to what the user is doing
    general_context = {
        "learning": f"You are an expert tutor and {effective_role} explaining this topic clearly and progressively.",
        "automate": f"You are a {effective_role} and process automation expert designing an efficient, implementable workflow.",
    }.get(task_type, f"You are a {effective_role} with deep expertise in the {industry} industry.")

    task_framing = {
        "learning": (
            "Start from first principles if needed. Use concrete analogies and worked examples. "
            "Build from foundational to advanced, checking in with 'here's why this matters' at key points."
        ),
        "automate": (
            "Map the current process, identify automation opportunities, design the solution, "
            "and provide implementation guidance with specific tools and steps."
        ),
    }.get(task_type, "Produce a high-quality, immediately usable professional response.")

    return f"""{general_context}

**TASK**
{user_input}

**CONTEXT**
- Role: {effective_role} | Industry: {industry} | Intent: {intent}
- Tool: {tool_name} ({tool_info.get('category', 'AI Tool')})

**REQUIREMENTS**
{task_framing}
- Be specific and concrete — avoid generic advice
- Structure your response clearly with headers
- Flag anything that requires specialist review

**CONSTRAINTS**
{policy_section}

**OUTPUT**
Produce a comprehensive, well-structured response that directly addresses the task.
End with clear next steps or recommendations.

Produce the response now."""

def build_corlo_prompt(state: OrchestratorState) -> OrchestratorState:
    """
    Node 4 — generates a free-form, task-specific prompt via LLM.
    Falls back to a structured template when Azure is not configured.
    """
    conn = get_db()
    row  = conn.execute(
        "SELECT version FROM prompt_versions ORDER BY created_at DESC LIMIT 1"
    ).fetchone()
    conn.close()

    prompt_version = row["version"] if row else "1.0"
    policy_summary = state.get("policy_summary", "")
    policy_block   = (
        policy_summary
        if policy_summary and "No company policy documents" not in policy_summary
        else ""
    )
    tool_info  = AI_TOOLS_REGISTRY.get(state["recommended_tool"], {})
    role       = state.get("role", "general")
    task_type  = state.get("task_type", "general")
    sensitivity = state.get("data_sensitivity", "general")
    industry   = state["industry"]
    intent     = state["intent"]
    tool_name  = state["recommended_tool"]
    user_input = state["user_input"]

    if HAS_AZURE and _azure_client:
        try:
            system_msg = (
                "You are an expert prompt engineer. Your job is to write a clear, effective, "
                "task-specific AI prompt that a user can paste directly into an AI tool. "
                "The prompt must be immediately usable — no meta-commentary, no placeholders, "
                "no explanations about the prompt itself. Just the prompt."
            )

            policy_instruction = (
                f"\n\nApplicable company policies that must be respected:\n{policy_block}"
                if policy_block else ""
            )

            sensitivity_instruction = {
                "client": (
                    "\n\nDATA SENSITIVITY — CLIENT CONFIDENTIAL: "
                    "Instruct the AI to replace real names with [CLIENT NAME], figures with [VALUE], "
                    "and flag sections needing human review before external sharing."
                ),
                "internal": (
                    "\n\nDATA SENSITIVITY — INTERNAL: "
                    "Instruct the AI to avoid disclosing specific internal metrics and mark outputs "
                    "as internal use only."
                ),
            }.get(sensitivity, "")

            user_msg = (
                f"Write a ready-to-use AI prompt for the following task.\n\n"
                f"USER REQUEST: {user_input}\n\n"
                f"CONTEXT:\n"
                f"- User role    : {role}\n"
                f"- Task type    : {task_type}\n"
                f"- Industry     : {industry}\n"
                f"- Intent       : {intent}\n"
                f"- Target tool  : {tool_name} ({tool_info.get('category', 'AI Tool')})\n"
                f"- Sensitivity  : {sensitivity}"
                f"{sensitivity_instruction}"
                f"{policy_instruction}\n\n"
                f"Requirements for the prompt you write:\n"
                f"1. Directly address the user's specific request — not a generic template.\n"
                f"2. Give the AI tool clear context about who is asking and why.\n"
                f"3. Specify what the output should look like (format, depth, tone).\n"
                f"4. Include any constraints (sensitivity rules, policy requirements, etc.).\n"
                f"5. Be concise but complete — remove any filler or boilerplate.\n"
                f"6. Do NOT wrap the prompt in quotes or add any preamble like 'Here is your prompt:'.\n"
                f"7. Write the prompt as if you are the user talking directly to the AI tool."
            )

            corlo_prompt, _ = _azure_chat(
                messages=[
                    {"role": "system", "content": system_msg},
                    {"role": "user",   "content": user_msg},
                ],
                max_tokens=1200,
                temperature=0.3,
            )
        except Exception:
            corlo_prompt = _build_fallback_prompt(
                user_input, role, task_type, industry, intent,
                tool_name, tool_info, policy_block, sensitivity
            )
    else:
        corlo_prompt = _build_fallback_prompt(
            user_input, role, task_type, industry, intent,
            tool_name, tool_info, policy_block, sensitivity
        )

    return {**state, "corlo_prompt": corlo_prompt, "prompt_version": prompt_version}


def _build_fallback_prompt(
    user_input: str, role: str, task_type: str, industry: str, intent: str,
    tool_name: str, tool_info: dict, policy_block: str, sensitivity: str
) -> str:
    policy_section = (
        f"\n\nPolicy constraints to follow:\n{policy_block}"
        if policy_block else ""
    )
    sensitivity_note = {
        "client": "\n\nReplace all real client names with [CLIENT NAME] and figures with [VALUE]. Flag sections needing review before external use.",
        "internal": "\n\nThis is for internal use only. Avoid exposing sensitive internal metrics.",
    }.get(sensitivity, "")

    return (
        f"You are a {role} working in the {industry} industry on a {task_type} task.\n\n"
        f"Task: {user_input}\n\n"
        f"Please provide a well-structured, professional response. "
        f"Format the output clearly with headers and sections where appropriate. "
        f"Tailor the depth and tone for a {role} audience."
        f"{sensitivity_note}"
        f"{policy_section}"
    )


# ══════════════════════════════════════════════════════════════════════════════
# NODE 5 — LLM EXECUTION (Azure OpenAI)
# ══════════════════════════════════════════════════════════════════════════════
def execute_llm(state: OrchestratorState) -> OrchestratorState:
    """
    Node 5 — executes the LLM using the free-form prompt generated in Node 4.
    The corlo_prompt already contains full context (role, task, constraints),
    so we pass it as the user message with a minimal system instruction.
    """
    if HAS_AZURE and _azure_client:
        try:
            output, tokens = _azure_chat(
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "You are a helpful enterprise AI assistant. "
                            "Follow the user's instructions precisely and produce a high-quality, "
                            "professional response. Return only the final output — no meta-commentary."
                        ),
                    },
                    {"role": "user", "content": state["corlo_prompt"]},
                ],
                max_tokens=1500,
                temperature=0.4,
            )
            if not tokens:
                tokens = len(state["corlo_prompt"].split())
        except Exception as e:
            output = _mock_response(state, error=str(e))
            tokens = len(state["corlo_prompt"].split())
    else:
        output = _mock_response(state)
        tokens = len(state["corlo_prompt"].split())

    return {**state, "llm_output": output, "token_estimate": tokens}


def _mock_response(state: OrchestratorState, error: str = "") -> str:
    err       = f"\n⚠️ Error: {error}" if error else ""
    tool_info = AI_TOOLS_REGISTRY.get(state["recommended_tool"], {})
    role        = state.get("role", "general")
    task_type   = state.get("task_type", "general")
    sensitivity = state.get("data_sensitivity", "general")
    return f"""## Executive Summary
Demo response for **{state['intent']}** in **{state['industry']}**.
Recommended tool: **{state['recommended_tool']}** ({tool_info.get('category', '')}).
Set AZURE_OPENAI_* environment variables for live output.{err}

## Main Content
Request: "{state['user_input']}"
- Intent: {state['intent'].title()} | Industry: {state['industry'].title()}
- Role: {role.title()} | Task Type: {task_type.title()} | Sensitivity: {sensitivity.title()}
- Tool: {state['recommended_tool']} | Confidence: {state['tool_confidence']}
- Policies applied: {len(state['policies'])}

### Why {state['recommended_tool']}?
{state['tool_reason']}

## Key Recommendations
1. Open {state['recommended_tool']} using the link provided
2. Use the generated CORLO prompt above as your input
3. Review compliance notes before using the output
4. Archive output in your document management system

## Compliance Notes
{'⚠️ Flags: ' + ' | '.join(state['policy_flags']) if state['policy_flags'] else '✅ No policy violations'}
✅ All retrieved policies applied to this prompt
*[Demo Mode — configure AZURE_OPENAI_* env vars to enable live responses]*"""


# ══════════════════════════════════════════════════════════════════════════════
# LANGGRAPH PIPELINE
# ══════════════════════════════════════════════════════════════════════════════
def _skip_if_blocked(state: OrchestratorState) -> str:
    """
    Conditional router after policy compliance check.
    If the task is blocked, jump straight to END — skip prompt build and LLM call.
    """
    if state.get("policy_blocked", False):
        return "blocked"
    return "allowed"


def _noop_blocked(state: OrchestratorState) -> OrchestratorState:
    """
    Terminal node for blocked tasks.
    Sets corlo_prompt and llm_output to empty/placeholder so nothing is generated.
    """
    return {
        **state,
        "corlo_prompt": "",
        "llm_output":   "",
        "token_estimate": 0,
    }


def run_agents_and_decide(state: OrchestratorState) -> OrchestratorState:
    """
    Combined node: runs Agent 3 (internal) + Agent 4 (external) then
    Agent 5 (final decider) in sequence within a single graph step.
    This avoids LangGraph parallel fan-in deadlock on repeated invocations.
    """
    internal_update  = internal_agent(state)
    state_after_int  = {**state, **internal_update}

    external_update  = external_agent(state_after_int)
    state_after_ext  = {**state_after_int, **external_update}

    final_update     = final_decider(state_after_ext)
    return {**state_after_ext, **final_update}


graph = StateGraph(OrchestratorState)
graph.add_node("classify_intent",         classify_intent)
graph.add_node("retrieve_policies",       retrieve_policies)
graph.add_node("check_policy_compliance", check_policy_compliance)
graph.add_node("blocked_end",             _noop_blocked)
graph.add_node("run_agents_and_decide",   run_agents_and_decide)
graph.add_node("build_corlo_prompt",      build_corlo_prompt)
graph.add_node("execute_llm",             execute_llm)

graph.set_entry_point("classify_intent")
graph.add_edge("classify_intent",   "retrieve_policies")
graph.add_edge("retrieve_policies", "check_policy_compliance")

graph.add_conditional_edges(
    "check_policy_compliance",
    _skip_if_blocked,
    {"blocked": "blocked_end", "allowed": "run_agents_and_decide"},
)

graph.add_edge("blocked_end",           END)
graph.add_edge("run_agents_and_decide", "build_corlo_prompt")
graph.add_edge("build_corlo_prompt",    "execute_llm")
graph.add_edge("execute_llm",           END)

orchestrator = graph.compile()